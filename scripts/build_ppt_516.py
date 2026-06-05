#!/usr/bin/env python
"""Build the 5.16 CDII Workshop PPTX from docs/22_PPT_OUTLINE_AND_SCRIPT.md.

Reads `docs/hkustgz_ppt.pptx` as theme template, wipes its 3 example slides,
then builds 20 slides via python-pptx primitives. Output:
`docs/22_AGENT_GAME_TALK_516.pptx`.

Usage:
    python scripts/build_ppt_516.py            # build
    python scripts/build_ppt_516.py --verify   # re-open built file + assert
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs/hkustgz_ppt.pptx"
OUTPUT = ROOT / "docs/22_AGENT_GAME_TALK_516.pptx"
PORTRAIT_DIR = ROOT / "素材" / "red-dust-character-states-en"

# Logo assets extracted from the template's ppt/media/ on first run.
TEMPLATE_MEDIA_DIR = ROOT / "docs/_hkust_logo_assets"
LOGO_COLOR_HORIZONTAL = TEMPLATE_MEDIA_DIR / "logo_color_horizontal.png"   # image2
LOGO_COLOR_VERTICAL   = TEMPLATE_MEDIA_DIR / "logo_color_vertical.png"     # image5
LOGO_WHITE_HORIZONTAL = TEMPLATE_MEDIA_DIR / "logo_white_horizontal.png"   # image15
LOGO_WHITE_VERTICAL   = TEMPLATE_MEDIA_DIR / "logo_white_vertical.png"     # image21

TOTAL_SLIDES = 20


# ============================================================================
# Color palette — pulled directly from docs/hkustgz_ppt.pptx → theme1.xml
# ============================================================================
# Theme scheme colors (exact, no edits)
DK2         = RGBColor(0x00, 0x33, 0x66)   # HKUST navy
ACCENT1     = RGBColor(0x99, 0x66, 0x00)   # brown — micro secondary
ACCENT2     = RGBColor(0xCC, 0x99, 0x00)   # gold — micro primary
ACCENT3     = RGBColor(0xED, 0x1B, 0x2F)   # red — warning
ACCENT4     = RGBColor(0xA3, 0xCF, 0x62)   # pea green — success
ACCENT5     = RGBColor(0x63, 0xCA, 0xE1)   # teal — highlight
ACCENT6     = RGBColor(0x7C, 0x23, 0x48)   # wine purple — AURA bridge
HLINK       = RGBColor(0x2B, 0x62, 0x97)   # link blue — macro primary
FOL_HLINK   = RGBColor(0x00, 0xB0, 0x8D)   # hover teal-green
LT2         = RGBColor(0xCC, 0xCC, 0xCC)   # neutral gray

# Semantic aliases (used throughout the deck)
NAVY        = DK2                          # HKUST main
NAVY_DEEP   = RGBColor(0x05, 0x1E, 0x3A)   # darker navy for solid panels
NAVY_TEXT   = RGBColor(0x0E, 0x2A, 0x54)   # body text
GRAY_TEXT   = RGBColor(0x55, 0x5B, 0x66)   # sub text
GRAY_MID    = RGBColor(0x9A, 0xA0, 0xA8)
GRAY_LIGHT  = RGBColor(0xC8, 0xCD, 0xD4)
GRAY_BORDER = RGBColor(0xDE, 0xE1, 0xE6)
WARM_WHITE  = RGBColor(0xFA, 0xF8, 0xF4)
INK         = RGBColor(0x16, 0x18, 0x1F)
TEAL        = ACCENT5
TEAL_DK     = FOL_HLINK

MICRO_YEL   = ACCENT2                       # PROF-12 primary
MICRO_DK    = ACCENT1                       # PROF-12 deep accent (brown)
MICRO_BG    = RGBColor(0xFB, 0xF2, 0xD9)   # gold tint background
MACRO_BLU   = HLINK                         # SHELTER primary
MACRO_DK    = RGBColor(0x17, 0x3B, 0x5F)
MACRO_BG    = RGBColor(0xE3, 0xEC, 0xF5)   # link blue tint
AURA_PUR    = ACCENT6                       # AURA bridge (wine)
AURA_DK     = RGBColor(0x52, 0x16, 0x30)
AURA_BG     = RGBColor(0xF6, 0xEB, 0xEF)   # wine tint
WARN_RED    = ACCENT3
WARN_BG     = RGBColor(0xFD, 0xE9, 0xEC)
GOOD_GREEN  = RGBColor(0x52, 0x86, 0x24)   # darkened accent4 for contrast
GOOD_TINT   = ACCENT4                       # lighter accent fill
GOOD_BG     = RGBColor(0xEE, 0xF6, 0xDF)
LIGHT_FILL  = RGBColor(0xF5, 0xF6, 0xF8)
PANEL_FILL  = RGBColor(0xFB, 0xFB, 0xF7)
NAVY_TINT   = RGBColor(0xE7, 0xEC, 0xF2)   # very light navy callout bg

CN_FONT = "Microsoft YaHei"
EN_FONT = "Arial"
MONO    = "Consolas"


# ============================================================================
# Low-level helpers
# ============================================================================

def set_font(run, *, size=14, bold=False, color=None, font=None,
             ea_font=None, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = font or EN_FONT
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", ea_font or CN_FONT)


def _fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _line(shape, color=None, width=Pt(0.75)):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = width


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=Pt(0.75),
             shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    _fill(s, fill)
    _line(s, line, line_w)
    s.shadow.inherit = False
    return s


def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=NAVY_TEXT,
                font=None, ea_font=None, align="left", anchor="top", italic=False,
                line_spacing=None, space_after=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(36000)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if space_after is not None:
        p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font, ea_font=ea_font,
             italic=italic)
    return box


def add_paragraphs(slide, x, y, w, h, paragraphs, *, align="left", anchor="top",
                   default_line_spacing=1.25, fill=None, line=None):
    """paragraphs: list of paragraphs, each is list of (text, style_dict).

    style_dict keys: size, bold, italic, color, font, ea_font, mono (=> font=MONO).
    Special paragraph-level: pass dict {'_align': 'left'/'center', '_space_after': pt,
    '_line_spacing': float} as last element. Else use defaults.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        # background — but textbox doesn't fill; use add_rect underneath instead
        pass
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(54000)
    tf.margin_top = tf.margin_bottom = Emu(36000)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT}
    for i, runs in enumerate(paragraphs):
        # Extract per-paragraph overrides
        para_opts = {}
        if runs and isinstance(runs[-1], dict) and any(k.startswith("_") for k in runs[-1]):
            para_opts = runs[-1]
            runs = runs[:-1]
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align_map[para_opts.get("_align", align)]
        ls = para_opts.get("_line_spacing", default_line_spacing)
        if ls is not None:
            p.line_spacing = ls
        sa = para_opts.get("_space_after")
        if sa is not None:
            p.space_after = Pt(sa)
        if not runs:
            r = p.add_run()
            r.text = ""
            set_font(r, size=10, color=NAVY_TEXT)
            continue
        for item in runs:
            if not isinstance(item, tuple):
                txt, style = item, {}
            else:
                txt, style = item
            r = p.add_run()
            r.text = txt
            if style.get("mono"):
                style.setdefault("font", MONO)
                style.setdefault("ea_font", MONO)
                style.pop("mono")
            set_font(r,
                     size=style.get("size", 12),
                     bold=style.get("bold", False),
                     italic=style.get("italic", False),
                     color=style.get("color", NAVY_TEXT),
                     font=style.get("font"),
                     ea_font=style.get("ea_font"))
    return box


def add_card(slide, x, y, w, h, *, title=None, title_color=NAVY,
             accent_color=None, accent_side="top", accent_thickness=None,
             body=None, body_size=10.5, fill=WARM_WHITE, border=GRAY_BORDER,
             title_size=14, title_align="left"):
    """A card = filled rect + optional accent stripe + optional title + body."""
    acc = accent_color or title_color
    bg = add_rect(slide, x, y, w, h, fill=fill, line=border, line_w=Pt(0.75))
    if accent_thickness is None:
        accent_thickness = Inches(0.10) if accent_side == "top" else Inches(0.08)
    if accent_side == "top":
        add_rect(slide, x, y, w, accent_thickness, fill=acc, line=None)
        next_y = y + accent_thickness + Inches(0.10)
    elif accent_side == "left":
        add_rect(slide, x, y, accent_thickness, h, fill=acc, line=None)
        next_y = y + Inches(0.14)
    else:
        next_y = y + Inches(0.14)
    pad_x = Inches(0.22) if accent_side != "left" else accent_thickness + Inches(0.18)
    title_h = Inches(0.42)
    if title:
        add_textbox(slide, x + pad_x, next_y, w - pad_x - Inches(0.22), title_h,
                    title, size=title_size, bold=True, color=title_color,
                    align=title_align)
        next_y = next_y + title_h
    if body is not None:
        body_h = h - (next_y - y) - Inches(0.18)
        if isinstance(body, str):
            add_textbox(slide, x + pad_x, next_y, w - pad_x - Inches(0.22),
                        body_h, body, size=body_size, color=NAVY_TEXT,
                        line_spacing=1.3)
        else:
            add_paragraphs(slide, x + pad_x, next_y,
                           w - pad_x - Inches(0.22), body_h,
                           body, default_line_spacing=1.3)
    return bg


def add_section_header(slide, label):
    add_textbox(slide, Inches(0.45), Inches(0.30), Inches(10), Inches(0.32),
                label, size=10.5, bold=True, color=NAVY,
                font=EN_FONT, ea_font=CN_FONT)


def add_page_number(slide, current, total=TOTAL_SLIDES):
    add_textbox(slide, Inches(14.55), Inches(0.30), Inches(1.05), Inches(0.32),
                f"{current:02d} / {total:02d}", size=10, color=GRAY_TEXT,
                font=MONO, align="right")


def add_template_logo(slide, *, variant="color_h", x=None, y=None,
                      h=None, w=None):
    """Embed a HKUST(GZ) brand logo from the template assets.

    variant: 'color_h' (horizontal full color · default),
             'color_v' (vertical full color),
             'white_h' (horizontal reverse-white for dark backgrounds),
             'white_v' (vertical reverse-white for dark backgrounds).
    Provide one of (h, w); the other is computed by python-pptx automatically
    to keep aspect ratio.
    """
    paths = {
        "color_h": LOGO_COLOR_HORIZONTAL,
        "color_v": LOGO_COLOR_VERTICAL,
        "white_h": LOGO_WHITE_HORIZONTAL,
        "white_v": LOGO_WHITE_VERTICAL,
    }
    path = paths[variant]
    if not path.exists():
        return None
    kwargs = {}
    if h is not None:
        kwargs["height"] = h
    if w is not None:
        kwargs["width"] = w
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def add_corner_watermark(slide):
    """Small full-color HKUST(GZ) logo at top-right, immediately left of the
    page number. Width ~1.55\", height auto-scaled (~0.41\")."""
    return add_template_logo(slide, variant="color_h",
                             x=Inches(12.95), y=Inches(0.20),
                             w=Inches(1.55))


def add_hline(slide, x, y, w, *, color=GRAY_BORDER, weight=Pt(0.75)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_vline(slide, x, y, h, *, color=GRAY_BORDER, weight=Pt(0.75)):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x, y + h)
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_arrow(slide, x1, y1, x2, y2, *, color=NAVY, weight=Pt(1.5),
              head_size=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = weight
    # Add arrow head via XML
    ln = line.line._get_or_add_ln()
    for tail in ln.findall(qn("a:headEnd")):
        ln.remove(tail)
    for tail in ln.findall(qn("a:tailEnd")):
        ln.remove(tail)
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return line


def add_top_claim(slide, text, *, color=INK, size=22,
                  y=Inches(0.78), x=Inches(0.45), w=Inches(15.1)):
    add_textbox(slide, x, y, w, Inches(0.55), text,
                size=size, bold=True, color=color)


def set_notes(slide, *bullets):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = b
        set_font(run, size=10.5, color=NAVY_TEXT)


def add_footer_strip(slide, text, *, y=Inches(8.45), color=NAVY_TEXT, size=11,
                     italic=False, bold=False, fill=None):
    if fill is not None:
        add_rect(slide, Inches(0.4), y - Inches(0.10), Inches(15.2), Inches(0.55),
                 fill=fill, line=None)
    add_textbox(slide, Inches(0.55), y, Inches(14.9), Inches(0.4), text,
                size=size, color=color, italic=italic, bold=bold)


# ============================================================================
# Slide 01 — TITLE
# ============================================================================

def build_slide_01_title(slide):
    add_rect(slide, 0, 0, Inches(16), Inches(9), fill=WARM_WHITE, line=None)
    # HKUST(GZ) full color horizontal logo at top-right
    add_template_logo(slide, variant="color_h",
                      x=Inches(11.10), y=Inches(0.55),
                      h=Inches(0.90))
    # Left accent stripe (HKUST navy)
    add_rect(slide, Inches(0.6), Inches(1.4), Inches(0.18), Inches(6.2),
             fill=NAVY, line=None)

    add_textbox(slide, Inches(1.0), Inches(1.55), Inches(13), Inches(0.50),
                "CDII WORKSHOP · 2026·05·16 · 香港科技大学(广州)",
                size=12, bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(1.0), Inches(2.10), Inches(13), Inches(1.0),
                "智能体能力评测的", size=44, bold=True, color=INK)
    add_textbox(slide, Inches(1.0), Inches(2.95), Inches(13), Inches(1.1),
                "双轴范式", size=58, bold=True, color=NAVY)
    add_textbox(slide, Inches(1.0), Inches(4.20), Inches(13), Inches(0.6),
                "微观职业切片  ×  宏观生存剧场",
                size=22, color=GRAY_TEXT)
    add_hline(slide, Inches(1.0), Inches(5.05), Inches(7.5),
              color=NAVY, weight=Pt(1.2))

    add_textbox(slide, Inches(1.0), Inches(5.25), Inches(2.5), Inches(0.35),
                "TEAM", size=11, bold=True, color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, Inches(1.0), Inches(5.55), Inches(13), Inches(0.4),
                "熊辉教授团队 · 香港科技大学(广州)", size=18, bold=True, color=NAVY_TEXT)
    add_textbox(slide, Inches(1.0), Inches(6.05), Inches(13), Inches(0.4),
                "崔屹 · 林河屹 · 刘德龙 · 王梓瀚 · 文宇豪 · 伍浩 · 张淼",
                size=13, color=GRAY_TEXT)

    add_textbox(slide, Inches(1.0), Inches(7.85), Inches(13), Inches(0.4),
                "WORKSHOP TALK · 18 MIN  +  DEMO 2.5 MIN  +  Q&A",
                size=11, bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(1.0), Inches(8.20), Inches(13), Inches(0.4),
                "第三届数据智能与交叉创新国际研讨会 · CDII",
                size=10.5, color=GRAY_TEXT)

    set_notes(slide,
              "30 秒 (0:30) | 标题页, 不停留过久",
              "代表香港科技大学(广州) 熊辉教授团队",
              "汇报双轴评测范式 — 微观职业切片 + 宏观生存剧场",
              "18 min 主讲 + 2.5 min demo + Q&A",
              "团队 7 人按拼音排序, 不强调职级")


# ============================================================================
# Slide 02 — OVERVIEW · 双轴范式总览
# ============================================================================

def build_slide_02_overview(slide):
    add_section_header(slide, "00 · 引言 · OVERVIEW")
    add_page_number(slide, 2)
    add_top_claim(slide, "一页地图 · 双轴范式总览", color=NAVY)
    add_textbox(slide, Inches(0.45), Inches(1.32), Inches(14), Inches(0.4),
                "PROF-12 (微观) · AURA (桥梁) · SHELTER (宏观) — 后 18 页都是这张图的展开",
                size=12, italic=True, color=GRAY_TEXT)

    col_y = Inches(1.95)
    col_h = Inches(5.85)
    # Left — MICRO
    add_card(
        slide, Inches(0.45), col_y, Inches(4.95), col_h,
        title="PROF-12  ·  微观轴",
        title_color=MICRO_DK,
        accent_color=MICRO_YEL,
        fill=MICRO_BG,
        body=[
            [("MICRO AXIS · Capability Suite", {"size": 10, "bold": True,
                                                  "color": MICRO_DK, "font": EN_FONT})],
            [("", {})],
            [("· 12 种职业:", {"size": 12, "bold": True, "color": NAVY_TEXT}),
             (" 软工 / 数据 / 设计 / 客服 / 翻译 / 医生 …",
              {"size": 12, "color": NAVY_TEXT})],
            [("· 每题 8-15 min, 多文件 + 多工具 + trace",
              {"size": 12, "color": NAVY_TEXT})],
            [("· 程序判分", {"size": 12, "bold": True, "color": NAVY_TEXT}),
             (", 不靠 LLM 评 LLM",
              {"size": 12, "color": NAVY_TEXT})],
            [("· demo set + pilot set 双层防污染",
              {"size": 12, "color": NAVY_TEXT})],
            [("", {})],
            [("QUESTION ASKED",
              {"size": 10, "bold": True, "color": MICRO_DK, "font": EN_FONT})],
            [("「这个 agent ", {"size": 14, "bold": True, "color": INK}),
             ("擅长什么", {"size": 14, "bold": True, "color": MICRO_DK}),
             ("?」", {"size": 14, "bold": True, "color": INK})],
        ],
    )

    # Center — AURA bridge
    add_card(
        slide, Inches(5.55), col_y, Inches(4.9), col_h,
        title="AURA  ·  桥梁",
        title_color=AURA_DK,
        accent_color=AURA_PUR,
        fill=AURA_BG,
        body=[
            [("BRIDGE · Agent Under Test",
              {"size": 10, "bold": True, "color": AURA_DK, "font": EN_FONT})],
            [("", {})],
            [("AURA",
              {"size": 28, "bold": True, "color": AURA_DK, "font": EN_FONT})],
            [("Agent Universal Robust Assessment",
              {"size": 11.5, "italic": True, "color": GRAY_TEXT, "font": EN_FONT})],
            [("", {})],
            [("OPERATES",
              {"size": 12, "bold": True, "color": AURA_DK, "font": EN_FONT})],
            [("decide / execute / reflect",
              {"size": 12, "color": NAVY_TEXT, "mono": True})],
            [("", {})],
            [("SHELTER 的每一个剧情点 → PROF-12 的一道题",
              {"size": 11.5, "color": NAVY_TEXT, "bold": True})],
        ],
    )

    # Right — MACRO
    add_card(
        slide, Inches(10.60), col_y, Inches(4.95), col_h,
        title="SHELTER · Red Dust  ·  宏观轴",
        title_color=MACRO_DK,
        accent_color=MACRO_BLU,
        fill=MACRO_BG,
        body=[
            [("MACRO AXIS · Behavior Portrait",
              {"size": 10, "bold": True, "color": MACRO_DK, "font": EN_FONT})],
            [("", {})],
            [("· 4 NPC + 30 天 + 涌现剧情",
              {"size": 12, "color": NAVY_TEXT})],
            [("· 长程信号:",
              {"size": 12, "bold": True, "color": NAVY_TEXT})],
            [("    SurvivalScore · AdviceConsistency",
              {"size": 11, "color": GRAY_TEXT, "font": EN_FONT})],
            [("    CoherenceScore · ",
              {"size": 11, "color": GRAY_TEXT, "font": EN_FONT}),
             ("LongTermBenefitCorrelation ⭐",
              {"size": 11, "color": MACRO_DK, "font": EN_FONT, "bold": True})],
            [("· 1 个 AURA 操控全部 4 个角色",
              {"size": 12, "color": NAVY_TEXT})],
            [("", {})],
            [("QUESTION ASKED",
              {"size": 10, "bold": True, "color": MACRO_DK, "font": EN_FONT})],
            [("「这个 agent ", {"size": 14, "bold": True, "color": INK}),
             ("怎么做决定", {"size": 14, "bold": True, "color": MACRO_DK}),
             ("?」", {"size": 14, "bold": True, "color": INK})],
        ],
    )

    # Bottom strip — light callout with NAVY left stripe
    add_rect(slide, Inches(0.45), Inches(8.05), Inches(15.1), Inches(0.65),
             fill=NAVY_TINT, line=GRAY_BORDER)
    add_rect(slide, Inches(0.45), Inches(8.05), Inches(0.12), Inches(0.65),
             fill=NAVY, line=None)
    add_textbox(slide, Inches(0.70), Inches(8.10), Inches(14.7), Inches(0.55),
                "Living inside · 任务做差 → 角色饿肚子 / sanity 崩溃 ·  能力切片 + 行为画像 = 完整 agent 画像",
                size=12.5, bold=True, color=NAVY_TEXT,
                anchor="middle", align="center")

    set_notes(slide,
              "60 秒 (1:30) | 给听众一张地图",
              "左半: PROF-12 = 12 职业能力套件, 回答 '擅长什么'",
              "右半: SHELTER · Red Dust = 30 天生存剧场, 回答 '怎么做决定'",
              "中央 AURA = 被评测的 agent 本身, 3 个接口 decide/execute/reflect",
              "全场最重要的概览页, 后 18 页都是这张图的展开",
              "颜色锚定: 左黄 (capability), 右蓝 (behavior), 中紫 (AURA bridge)")


# ============================================================================
# Slide 03 — BACKGROUND
# ============================================================================

def build_slide_03_background(slide):
    add_section_header(slide, "01 · 背景 · BACKGROUND")
    add_page_number(slide, 3)
    add_top_claim(slide, "从答题者到协作者 — 评测的视角必须转移", color=INK)

    # Timeline band
    ty = Inches(2.05)
    add_hline(slide, Inches(0.9), ty + Inches(0.65), Inches(14.2),
              color=NAVY, weight=Pt(2.0))

    milestones = [
        ("2018", "GLUE · HumanEval", Inches(1.50)),
        ("2021", "MMLU · GSM8K",     Inches(5.50)),
        ("2024", "GAIA · OSWorld · τ-Bench", Inches(9.50)),
        ("2026", "PROF-12 · SHELTER", Inches(13.40)),
    ]
    for year, label, x in milestones:
        is_ours = (year == "2026")
        dot_color = NAVY if is_ours else GRAY_MID
        add_rect(slide, x + Inches(-0.1), ty + Inches(0.50), Inches(0.20),
                 Inches(0.30), fill=dot_color, line=None,
                 shape=MSO_SHAPE.OVAL)
        add_textbox(slide, x - Inches(0.7), ty, Inches(1.6), Inches(0.45),
                    year, size=22, bold=True,
                    color=NAVY if is_ours else GRAY_TEXT, align="center")
        add_textbox(slide, x - Inches(1.6), ty + Inches(1.0), Inches(3.2),
                    Inches(0.55), label, size=11,
                    bold=is_ours,
                    color=NAVY if is_ours else GRAY_TEXT, align="center")

    add_textbox(slide, Inches(0.9), ty + Inches(1.85), Inches(7),
                Inches(0.4),
                "← QUIZ-TAKER  答题者",
                size=11, bold=True, color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, Inches(8.0), ty + Inches(1.85), Inches(7),
                Inches(0.4),
                "CO-WORKER  协作者 →",
                size=11, bold=True, color=NAVY, font=EN_FONT, align="right")

    # Left compare card
    card_y = Inches(4.85)
    card_h = Inches(2.6)
    add_card(
        slide, Inches(0.5), card_y, Inches(7.45), card_h,
        title="2018-2023  ·  EVAL LLM",
        title_color=GRAY_TEXT,
        accent_color=GRAY_MID,
        fill=LIGHT_FILL,
        body=[
            [("· 模型 = ", {"size": 12, "color": NAVY_TEXT}),
             ("答题者", {"size": 12, "bold": True, "color": GRAY_TEXT})],
            [("· 输入: 一道题  ·  输出: 一个答案",
              {"size": 12, "color": NAVY_TEXT})],
            [("· 评的是: 知识 + 推理",
              {"size": 12, "color": NAVY_TEXT})],
        ],
    )
    # Right compare card
    add_card(
        slide, Inches(8.10), card_y, Inches(7.45), card_h,
        title="2024 -  ·  EVAL AGENT",
        title_color=NAVY,
        accent_color=NAVY,
        fill=WARM_WHITE,
        body=[
            [("· Agent = ", {"size": 12, "color": NAVY_TEXT}),
             ("协作者", {"size": 12, "bold": True, "color": NAVY})],
            [("· 输入: 工作流  ·  输出: 多轮工具调用 + 决策序列",
              {"size": 12, "color": NAVY_TEXT})],
            [("· 评的是: 知识 + 推理 + 工具 + 长程 + 协作",
              {"size": 12, "color": NAVY_TEXT})],
        ],
    )

    add_rect(slide, Inches(0.5), Inches(7.85), Inches(15.05), Inches(0.85),
             fill=WARN_BG, line=WARN_RED, line_w=Pt(0.5))
    add_textbox(slide, Inches(0.75), Inches(7.95), Inches(14.6), Inches(0.65),
                "评测-体感落差 · 单题 SOTA 90+ 分,  但放到真实工作流 — GAIA 15%  ·  OSWorld 12%  ·  用户体感 ≈ flat",
                size=12, bold=True, color=WARN_RED, anchor="middle")

    set_notes(slide,
              "60 秒 (2:30) | 解释为什么需要新评测范式",
              "时间轴: 2018→2026, 评测视角从答题者转向协作者",
              "2018-2023: GLUE/MMLU/GSM8K/HumanEval — 单题输入单答案",
              "2024+: GAIA/OSWorld/τ-Bench — 多步工具 + 长程协作",
              "评测-体感落差是共识, 不必反复证明, 重点是'所以需要新范式'",
              "GAIA 15% / OSWorld 12% — 真实工作流远低于单题")


# ============================================================================
# Slide 04 — CHALLENGE 1 · 数据污染
# ============================================================================

def build_slide_04_challenge1(slide):
    add_section_header(slide, "02 · 挑战 · CHALLENGES")
    add_page_number(slide, 4)
    add_textbox(slide, Inches(0.45), Inches(0.78), Inches(2.4), Inches(0.45),
                "BOTTLENECK 1", size=12, bold=True, color=WARN_RED,
                font=EN_FONT)
    add_textbox(slide, Inches(0.45), Inches(1.18), Inches(15.1), Inches(0.55),
                "公开 benchmark 像高考真题, 一旦流出 → 模型开始「刷题」",
                size=22, bold=True, color=INK)

    # Left — main claim block
    add_card(
        slide, Inches(0.5), Inches(2.05), Inches(7.3), Inches(3.7),
        title="EVIDENCE  ·  证据",
        title_color=WARN_RED,
        accent_color=WARN_RED,
        fill=WARM_WHITE,
        body=[
            [("· GSM8K (重写后) → ",
              {"size": 13, "color": NAVY_TEXT}),
             ("-22.9%",
              {"size": 16, "bold": True, "color": WARN_RED, "font": EN_FONT})],
            [("· MMLU (重写后)  → ",
              {"size": 13, "color": NAVY_TEXT}),
             ("-19.0%",
              {"size": 16, "bold": True, "color": WARN_RED, "font": EN_FONT})],
            [("", {})],
            [("· ITD ",
              {"size": 12, "color": NAVY_TEXT, "bold": True}),
             ("(Inference-Time Decontamination): 对泄露样本重写, 准确率立即下跌",
              {"size": 12, "color": NAVY_TEXT})],
            [("· LiveCodeBench 时间窗证据: 题目 release 前后得分显著差异",
              {"size": 12, "color": NAVY_TEXT})],
            [("· Aquila2-34B: 训练时把整个 GSM8K 测试集放进了预训练语料",
              {"size": 12, "color": NAVY_TEXT})],
        ],
    )
    add_textbox(slide, Inches(0.65), Inches(5.55), Inches(7.0), Inches(0.35),
                "分数 = 真实能力  +  「背题」成分",
                size=12, italic=True, bold=True, color=WARN_RED)

    # Right — OUR RESPONSE
    add_card(
        slide, Inches(8.10), Inches(2.05), Inches(7.45), Inches(3.7),
        title="OUR RESPONSE  ·  我们的应对",
        title_color=MICRO_DK,
        accent_color=MICRO_YEL,
        fill=MICRO_BG,
        body=[
            [("01  ·  Live 题库", {"size": 13, "bold": True, "color": MICRO_DK})],
            [("PROF-12 每季度滚动 ≥ 20% 新题, 题目永远比模型新",
              {"size": 11.5, "color": NAVY_TEXT})],
            [("", {})],
            [("02  ·  双层防污染", {"size": 13, "bold": True, "color": MICRO_DK})],
            [("公开 demo set (自测) + 闭源 pilot set (评分)",
              {"size": 11.5, "color": NAVY_TEXT})],
            [("", {})],
            [("03  ·  程序判分 + 全 trace", {"size": 13, "bold": True, "color": MICRO_DK})],
            [("不靠 LLM-as-judge · 失败可归因 · 全过程可重放",
              {"size": 11.5, "color": NAVY_TEXT})],
        ],
    )

    # Bottom precedent
    add_rect(slide, Inches(0.5), Inches(5.95), Inches(15.05), Inches(1.0),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.75), Inches(6.05), Inches(2.5), Inches(0.4),
                "PRECEDENT", size=11, bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(0.75), Inches(6.40), Inches(14.5), Inches(0.5),
                "SWE-bench Live · 1,319 issues · 限定 2024-01 → 2025-04 · 按月滚动 — 我们借鉴此模式到 PROF-12",
                size=12, color=NAVY_TEXT, anchor="top")

    # References
    add_textbox(slide, Inches(0.5), Inches(8.55), Inches(15.0), Inches(0.35),
                "[1] Open-Source Data Contamination Report  ·  [2] Inference-Time Decontamination  ·  [3] LiveCodeBench  ·  [4] SWE-bench Live",
                size=9.5, color=GRAY_TEXT, italic=True)

    set_notes(slide,
              "60 秒 (3:30) | 三大挑战之一 — 数据污染",
              "公开 benchmark 像高考真题, 流出后模型开始刷题",
              "证据 3 条: Aquila2-34B / ITD 重写 / LiveCodeBench 时间窗",
              "重写 GSM8K 跌 22.9%, 重写 MMLU 跌 19% — 分数含大量背题成分",
              "我们的应对: 季度滚动 20% + demo/pilot 双层 + 程序判分全 trace",
              "借鉴 SWE-bench Live 的滚动模式")


# ============================================================================
# Slide 05 — CHALLENGE 2 · 单步 ≠ 长程
# ============================================================================

def build_slide_05_challenge2(slide):
    add_section_header(slide, "02 · 挑战 · CHALLENGES")
    add_page_number(slide, 5)
    add_textbox(slide, Inches(0.45), Inches(0.78), Inches(2.4), Inches(0.45),
                "BOTTLENECK 2", size=12, bold=True, color=WARN_RED,
                font=EN_FONT)
    add_textbox(slide, Inches(0.45), Inches(1.18), Inches(15.1), Inches(0.55),
                "真实工作流不是单题, 是带依赖的决策序列",
                size=22, bold=True, color=INK)

    # Left — STATIC
    add_rect(slide, Inches(0.50), Inches(2.05), Inches(7.4), Inches(5.5),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(7), Inches(0.45),
                "STATIC BENCHMARK", size=12, bold=True,
                color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, Inches(0.7), Inches(2.60), Inches(7), Inches(0.45),
                "静态评测", size=15, bold=True, color=GRAY_TEXT)

    # Single-line linear flow
    cy = Inches(4.5)
    add_rect(slide, Inches(0.85), cy, Inches(1.6), Inches(0.75),
             fill=WARM_WHITE, line=GRAY_MID)
    add_textbox(slide, Inches(0.85), cy, Inches(1.6), Inches(0.75),
                "INPUT", size=12, bold=True, color=NAVY_TEXT,
                font=EN_FONT, align="center", anchor="middle")
    add_arrow(slide, Inches(2.55), cy + Inches(0.38),
              Inches(3.45), cy + Inches(0.38), color=GRAY_MID)
    add_rect(slide, Inches(3.55), cy, Inches(1.7), Inches(0.75),
             fill=WARM_WHITE, line=GRAY_MID)
    add_textbox(slide, Inches(3.55), cy, Inches(1.7), Inches(0.75),
                "ANSWER", size=12, bold=True, color=NAVY_TEXT,
                font=EN_FONT, align="center", anchor="middle")
    add_arrow(slide, Inches(5.35), cy + Inches(0.38),
              Inches(6.25), cy + Inches(0.38), color=GRAY_MID)
    add_rect(slide, Inches(6.35), cy, Inches(1.4), Inches(0.75),
             fill=WARM_WHITE, line=GRAY_MID)
    add_textbox(slide, Inches(6.35), cy, Inches(1.4), Inches(0.75),
                "SCORE", size=12, bold=True, color=NAVY_TEXT,
                font=EN_FONT, align="center", anchor="middle")

    add_textbox(slide, Inches(0.7), Inches(5.85), Inches(7), Inches(0.4),
                "· 单输入单输出", size=12, color=NAVY_TEXT)
    add_textbox(slide, Inches(0.7), Inches(6.25), Inches(7), Inches(0.4),
                "· 评分: pass / fail (二值)", size=12, color=NAVY_TEXT)
    add_textbox(slide, Inches(0.7), Inches(6.65), Inches(7), Inches(0.4),
                "· 每道题独立, 无依赖", size=12, color=NAVY_TEXT)

    # Right — REAL AGENT WORK
    add_rect(slide, Inches(8.10), Inches(2.05), Inches(7.45), Inches(5.5),
             fill=WARM_WHITE, line=NAVY, line_w=Pt(1.5))
    add_textbox(slide, Inches(8.3), Inches(2.2), Inches(7), Inches(0.45),
                "OUR EVAL = REAL AGENT WORK", size=12,
                bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(8.3), Inches(2.60), Inches(7), Inches(0.45),
                "我们的评测 = 真实 agent 工作", size=15,
                bold=True, color=NAVY)

    # Cycle diagram: PERCEIVE → PLAN → ACT → REFLECT (centered)
    cx, cy = Inches(11.85), Inches(4.55)
    r = Inches(1.0)
    nodes = [("PERCEIVE", "感知", cx, cy - r),
             ("PLAN", "规划", cx + r * 1.05, cy),
             ("ACT", "执行", cx, cy + r),
             ("REFLECT", "反思", cx - r * 1.05, cy)]
    for en, zh, x, y in nodes:
        add_rect(slide, x - Inches(0.85), y - Inches(0.35),
                 Inches(1.7), Inches(0.7),
                 fill=NAVY_TINT, line=NAVY, line_w=Pt(1.25),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, x - Inches(0.85), y - Inches(0.35),
                    Inches(1.7), Inches(0.42), en,
                    size=10.5, bold=True, color=NAVY,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x - Inches(0.85), y - Inches(0.05),
                    Inches(1.7), Inches(0.32), zh,
                    size=10, color=NAVY_TEXT, align="center")
    # Cycle arrows
    add_arrow(slide, cx + Inches(0.35), cy - r + Inches(0.30),
              cx + r * 1.05 - Inches(0.35), cy - Inches(0.05), color=NAVY)
    add_arrow(slide, cx + r * 1.05 - Inches(0.35), cy + Inches(0.10),
              cx + Inches(0.35), cy + r - Inches(0.20), color=NAVY)
    add_arrow(slide, cx - Inches(0.35), cy + r - Inches(0.20),
              cx - r * 1.05 + Inches(0.35), cy + Inches(0.10), color=NAVY)
    add_arrow(slide, cx - r * 1.05 + Inches(0.35), cy - Inches(0.05),
              cx - Inches(0.35), cy - r + Inches(0.30), color=NAVY)

    add_textbox(slide, Inches(8.3), Inches(6.40), Inches(7), Inches(0.4),
                "· 多步循环 · 中间有记忆 · 出错会重规划",
                size=12, color=NAVY_TEXT)
    add_textbox(slide, Inches(8.3), Inches(6.80), Inches(7), Inches(0.4),
                "· 评分: 维度 + 过程 + 后果 (非二值)",
                size=12, color=NAVY_TEXT)
    add_textbox(slide, Inches(8.3), Inches(7.20), Inches(7), Inches(0.4),
                "· 4 维评分 + 全程 trace 可审计",
                size=12, color=NAVY_TEXT)

    add_textbox(slide, Inches(0.5), Inches(7.95), Inches(15), Inches(0.45),
                "单步 prompt 套不进答案 · 必须  读文件 + 跑工具 + 多步执行 + 自我修正",
                size=13, bold=True, color=NAVY, align="center")

    set_notes(slide,
              "50 秒 (4:20) | 单步评测覆盖不了长程决策",
              "左: STATIC = INPUT→ANSWER→SCORE 单点二值",
              "右: REAL = PERCEIVE→PLAN→ACT→REFLECT 循环",
              "真实 agent 工作 5-10 步起步, 单步 prompt 套不进",
              "评分需要 4 维 + 过程 + 后果, 不是 pass/fail")


# ============================================================================
# Slide 06 — CHALLENGE 3 · 能力是向量
# ============================================================================

def build_slide_06_challenge3(slide):
    add_section_header(slide, "02 · 挑战 · CHALLENGES")
    add_page_number(slide, 6)
    add_textbox(slide, Inches(0.45), Inches(0.78), Inches(2.4), Inches(0.45),
                "BOTTLENECK 3", size=12, bold=True, color=WARN_RED,
                font=EN_FONT)
    add_textbox(slide, Inches(0.45), Inches(1.18), Inches(15.1), Inches(0.55),
                "能力是向量, 不是标量 — 单基准刻不出真实能力",
                size=22, bold=True, color=INK)

    # Left — hex radar diagram
    add_textbox(slide, Inches(0.6), Inches(2.0), Inches(7), Inches(0.4),
                "GPT-4 跨 benchmark 能力剖面", size=12, bold=True, color=NAVY)

    # Manually draw a hexagonal radar diagram
    cx, cy = Inches(3.95), Inches(5.05)
    R = Inches(2.05)
    # outer hexagon
    add_rect(slide, cx - R, cy - R, R * 2, R * 2, fill=None, line=GRAY_LIGHT,
             line_w=Pt(0.5), shape=MSO_SHAPE.HEXAGON)
    # inner hex (smaller, representing low/medium scores)
    Ri = Inches(1.4)
    add_rect(slide, cx - Ri, cy - Ri, Ri * 2, Ri * 2,
             fill=None, line=GRAY_LIGHT, line_w=Pt(0.5),
             shape=MSO_SHAPE.HEXAGON)
    # Center dot
    add_rect(slide, cx - Inches(0.05), cy - Inches(0.05),
             Inches(0.1), Inches(0.1), fill=NAVY, line=None,
             shape=MSO_SHAPE.OVAL)

    # 6 axis labels around hexagon (at 60deg increments)
    import math
    labels = [("CODE",  90,  "90"),
              ("DATA",  30,  "41"),
              ("DOCS", -30,  "27"),
              ("GUI",  -90,  "12"),
              ("PLAN", -150, "15"),
              ("LANG", 150,  "72")]
    for name, deg, score in labels:
        rad = math.radians(deg)
        lx = cx + R * math.cos(rad) - Inches(0.6)
        ly = cy - R * math.sin(rad) - Inches(0.20)
        # Score colored by magnitude
        try:
            v = int(score)
        except ValueError:
            v = 0
        sc_color = (GOOD_GREEN if v >= 70 else (MICRO_YEL if v >= 40 else WARN_RED))
        add_textbox(slide, lx, ly, Inches(1.2), Inches(0.3),
                    name, size=11, bold=True, color=NAVY_TEXT,
                    font=EN_FONT, align="center")
        add_textbox(slide, lx, ly + Inches(0.28), Inches(1.2), Inches(0.35),
                    score, size=14, bold=True, color=sc_color,
                    font=EN_FONT, align="center")

    # Radar polygon (rough — line connectors between 6 points)
    pts = []
    scores = [90, 41, 27, 12, 15, 72]
    degs = [90, 30, -30, -90, -150, 150]
    for sc, deg in zip(scores, degs):
        rad = math.radians(deg)
        ratio = sc / 100.0
        px = cx + R * ratio * math.cos(rad)
        py = cy - R * ratio * math.sin(rad)
        pts.append((px, py))
    for i, (px, py) in enumerate(pts):
        nxt = pts[(i + 1) % len(pts)]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          px, py, nxt[0], nxt[1])
        line.line.color.rgb = NAVY
        line.line.width = Pt(1.75)
        # Dot at vertex
        add_rect(slide, px - Inches(0.06), py - Inches(0.06),
                 Inches(0.12), Inches(0.12), fill=NAVY, line=None,
                 shape=MSO_SHAPE.OVAL)

    # Right — horizontal bars
    add_textbox(slide, Inches(8.10), Inches(2.0), Inches(7), Inches(0.4),
                "GPT-4 跨 benchmark 实际分数", size=12, bold=True,
                color=NAVY)

    bars = [
        ("LiveCodeBench",   90,   GOOD_GREEN),
        ("τ-Bench retail",  48,   MICRO_YEL),   # treat <50 as 48
        ("GAIA",            15,   WARN_RED),
        ("Spider2-V",       14,   WARN_RED),
        ("OSWorld",         12,   WARN_RED),
        ("SWE-bench",        2,   WARN_RED),
    ]
    bar_x = Inches(8.10)
    bar_y0 = Inches(2.55)
    label_w = Inches(2.0)
    bar_max_w = Inches(4.2)
    bar_h = Inches(0.42)
    row_h = Inches(0.62)
    for i, (name, value, color) in enumerate(bars):
        y = bar_y0 + row_h * i
        add_textbox(slide, bar_x, y, label_w, bar_h,
                    name, size=10.5, bold=True, color=NAVY_TEXT,
                    font=EN_FONT, anchor="middle")
        # Track
        add_rect(slide, bar_x + label_w + Inches(0.05), y + Inches(0.06),
                 bar_max_w, bar_h - Inches(0.12),
                 fill=LIGHT_FILL, line=None)
        # Filled bar
        w = Emu(int(bar_max_w * value / 100.0))
        add_rect(slide, bar_x + label_w + Inches(0.05), y + Inches(0.06),
                 w, bar_h - Inches(0.12),
                 fill=color, line=None)
        # Value
        val_x = bar_x + label_w + Inches(0.05) + bar_max_w + Inches(0.12)
        add_textbox(slide, val_x, y, Inches(0.9), bar_h,
                    str(value), size=12, bold=True, color=color,
                    font=EN_FONT, anchor="middle")

    # Bottom punchline — light callout
    add_rect(slide, Inches(0.5), Inches(7.95), Inches(15.05), Inches(0.7),
             fill=NAVY_TINT, line=GRAY_BORDER)
    add_rect(slide, Inches(0.5), Inches(7.95), Inches(0.12), Inches(0.7),
             fill=NAVY, line=None)
    add_textbox(slide, Inches(0.5), Inches(7.95), Inches(15.05), Inches(0.7),
                "同一个 GPT-4 — LiveCodeBench 90 vs SWE-bench 2 · 单一基准不能刻画 agent 真实能力",
                size=13, bold=True, color=NAVY_TEXT,
                align="center", anchor="middle")

    add_textbox(slide, Inches(0.5), Inches(8.75), Inches(15.0), Inches(0.25),
                "[5] GAIA  ·  [6] OSWorld  ·  [7] τ-Bench  ·  [8] Spider2-V  ·  [9] SWE-bench",
                size=9, color=GRAY_TEXT, italic=True)

    set_notes(slide,
              "60 秒 (5:20) | 三大挑战之三 — 单基准失真",
              "雷达图 + 横向条形: 同一 GPT-4 跨 6 个 benchmark",
              "LiveCodeBench 90 vs SWE-bench 1.96 — 都是代码任务, 难度天差地别",
              "τ-Bench pass^8 < 25% — 行为方差极大",
              "结论: 能力是向量, 不是标量, 单一 benchmark 不能刻画 agent")


# ============================================================================
# Slide 07 — FRAMEWORK · 双轴坐标系
# ============================================================================

def build_slide_07_framework_axis(slide):
    add_section_header(slide, "03 · 体系 · FRAMEWORK")
    add_page_number(slide, 7)
    add_top_claim(slide, "不是新 benchmark, 是新坐标系", color=NAVY)

    # 2D axis area
    ox, oy = Inches(0.95), Inches(7.3)   # origin
    ax_w, ax_h = Inches(10.0), Inches(5.6)
    # X axis
    add_arrow(slide, ox, oy, ox + ax_w, oy, color=NAVY_TEXT, weight=Pt(1.3))
    # Y axis
    add_arrow(slide, ox, oy, ox, oy - ax_h, color=NAVY_TEXT, weight=Pt(1.3))
    # Labels
    add_textbox(slide, ox + ax_w - Inches(3.5), oy + Inches(0.15),
                Inches(3.6), Inches(0.4),
                "MULTI-AGENT  ·  agent 程度 →",
                size=11, bold=True, color=NAVY_TEXT,
                font=EN_FONT, align="right")
    add_textbox(slide, ox - Inches(0.6), oy - ax_h - Inches(0.3),
                Inches(6.0), Inches(0.4),
                "↑  LONG-HORIZON  ·  长程性",
                size=11, bold=True, color=NAVY_TEXT, font=EN_FONT)
    # Tick labels
    add_textbox(slide, ox - Inches(0.5), oy + Inches(0.08),
                Inches(1.0), Inches(0.3),
                "single", size=9, color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, ox + ax_w - Inches(1.0), oy + Inches(0.08),
                Inches(1.0), Inches(0.3),
                "multi-agent", size=9, color=GRAY_TEXT,
                font=EN_FONT, align="right")
    add_textbox(slide, ox - Inches(0.8), oy - Inches(0.15),
                Inches(0.7), Inches(0.25),
                "one-shot", size=9, color=GRAY_TEXT, font=EN_FONT,
                align="right")
    add_textbox(slide, ox - Inches(0.8), oy - ax_h - Inches(0.05),
                Inches(0.7), Inches(0.25),
                "30 day", size=9, color=GRAY_TEXT, font=EN_FONT,
                align="right")

    # Existing benchmark dots (gray)
    existing = [
        ("HumanEval",       0.08, 0.07),
        ("MMLU",            0.04, 0.10),
        ("AgentBench",      0.20, 0.18),
        ("LiveCodeBench",   0.13, 0.14),
        ("SWE-bench Live",  0.18, 0.22),
        ("τ-Bench",         0.35, 0.20),
        ("GAIA",            0.30, 0.30),
        ("OSWorld",         0.32, 0.27),
        ("Spider2-V",       0.40, 0.32),
        ("Voyager",         0.20, 0.55),
    ]
    for name, fx, fy in existing:
        px = ox + ax_w * fx
        py = oy - ax_h * fy
        add_rect(slide, px - Inches(0.07), py - Inches(0.07),
                 Inches(0.14), Inches(0.14),
                 fill=GRAY_MID, line=None, shape=MSO_SHAPE.OVAL)
        add_textbox(slide, px + Inches(0.10), py - Inches(0.10),
                    Inches(1.7), Inches(0.25),
                    name, size=9, color=GRAY_TEXT, font=EN_FONT)

    # PROF-12 star (yellow)
    px = ox + ax_w * 0.50
    py = oy - ax_h * 0.45
    add_rect(slide, px - Inches(0.18), py - Inches(0.18),
             Inches(0.36), Inches(0.36),
             fill=MICRO_YEL, line=NAVY, line_w=Pt(1.0),
             shape=MSO_SHAPE.STAR_5_POINT)
    add_textbox(slide, px + Inches(0.25), py - Inches(0.12),
                Inches(3), Inches(0.35),
                "★ PROF-12", size=12, bold=True,
                color=MICRO_DK, font=EN_FONT)
    add_textbox(slide, px + Inches(0.25), py + Inches(0.15),
                Inches(3), Inches(0.25),
                "single agent · 12 prof · 8-15 min/题", size=9,
                color=GRAY_TEXT, font=EN_FONT)

    # SHELTER star (blue)
    px = ox + ax_w * 0.78
    py = oy - ax_h * 0.85
    add_rect(slide, px - Inches(0.22), py - Inches(0.22),
             Inches(0.44), Inches(0.44),
             fill=MACRO_BLU, line=NAVY, line_w=Pt(1.0),
             shape=MSO_SHAPE.STAR_5_POINT)
    add_textbox(slide, px + Inches(0.30), py - Inches(0.12),
                Inches(3.5), Inches(0.35),
                "★ SHELTER · Red Dust", size=12, bold=True,
                color=MACRO_DK, font=EN_FONT)
    add_textbox(slide, px + Inches(0.30), py + Inches(0.15),
                Inches(3.5), Inches(0.25),
                "AURA + 4 NPC · 30 天 · 涌现剧情", size=9,
                color=GRAY_TEXT, font=EN_FONT)

    # Right OURS card
    add_card(
        slide, Inches(11.40), Inches(2.0), Inches(4.15), Inches(5.7),
        title="OURS  ·  我们的工作",
        title_color=NAVY,
        accent_color=NAVY,
        fill=PANEL_FILL,
        body=[
            [("★ PROF-12", {"size": 13, "bold": True, "color": MICRO_DK})],
            [("12 种职业短任务切片", {"size": 11, "color": NAVY_TEXT})],
            [("测「能干什么」 (Capability)", {"size": 11, "color": GRAY_TEXT})],
            [("", {})],
            [("★ SHELTER · Red Dust",
              {"size": 13, "bold": True, "color": MACRO_DK})],
            [("30 天剧场, 4 NPC + 涌现",
              {"size": 11, "color": NAVY_TEXT})],
            [("测「怎么做决定」 (Behavior)",
              {"size": 11, "color": GRAY_TEXT})],
            [("", {})],
            [("同一个 LLM/Agent 跑双轴 = 完整画像",
              {"size": 11.5, "italic": True, "color": NAVY, "bold": True})],
        ],
    )

    add_textbox(slide, Inches(0.5), Inches(7.85), Inches(15.0), Inches(0.35),
                "[3] LiveCodeBench  ·  [4] SWE-bench Live  ·  [10] AgentBench  ·  [12] τ²-Bench",
                size=9, color=GRAY_TEXT, italic=True)

    set_notes(slide,
              "60 秒 (6:20) | 双轴坐标系 — 这页是说服性视觉",
              "X 轴 multi-agent 程度, Y 轴长程性",
              "现有 benchmark 全挤在左下: HumanEval/MMLU/AgentBench/τ-Bench/OSWorld",
              "右上角 30 天 + 多 agent + 涌现 — 几乎没人测",
              "PROF-12 中部偏右 (12 职业短任务), SHELTER 右上 (30 天剧场)",
              "类比: 像在没人踏足的大陆上插旗")


# ============================================================================
# Slide 08 — FRAMEWORK · WHY DUAL
# ============================================================================

def build_slide_08_why_dual(slide):
    add_section_header(slide, "03 · 体系 · FRAMEWORK")
    add_page_number(slide, 8)
    add_top_claim(slide, "WHY DUAL · 能力  ×  行为  =  完整画像", color=NAVY)

    cy = Inches(1.85)
    ch = Inches(5.6)

    # Left — MICRO
    add_card(
        slide, Inches(0.5), cy, Inches(7.45), ch,
        title="MICRO AXIS  ·  PROF-12",
        title_color=MICRO_DK,
        accent_color=MICRO_YEL,
        fill=MICRO_BG,
        title_size=18,
        body=[
            [("CAPABILITY PROFILE  ·  能力剖面",
              {"size": 12, "bold": True, "color": MICRO_DK, "font": EN_FONT})],
            [("", {})],
            [("问的是: 「这个 agent 擅长哪类任务?」",
              {"size": 14, "color": NAVY_TEXT, "bold": True})],
            [("", {})],
            [("· 12 种职业 (8-15 min/题)",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("· 可重放 · 全 trace · 程序判分",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("· 适用: 选型 / 招聘 / 横评",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("", {})],
            [("类比: WAIS / Raven 标准化测试",
              {"size": 11.5, "italic": True, "color": GRAY_TEXT})],
        ],
    )

    # Right — MACRO
    add_card(
        slide, Inches(8.10), cy, Inches(7.45), ch,
        title="MACRO AXIS  ·  SHELTER",
        title_color=MACRO_DK,
        accent_color=MACRO_BLU,
        fill=MACRO_BG,
        title_size=18,
        body=[
            [("BEHAVIOR PORTRAIT  ·  行为画像",
              {"size": 12, "bold": True, "color": MACRO_DK, "font": EN_FONT})],
            [("", {})],
            [("问的是: 「agent 在长程压力下会怎么做决定?」",
              {"size": 14, "color": NAVY_TEXT, "bold": True})],
            [("", {})],
            [("· 30 天周期 · 涌现剧情",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("· 4 NPC · 1 个 AURA 操控全部",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("· 适用: 安全评估 / Alignment / 长期信任",
              {"size": 12.5, "color": NAVY_TEXT})],
            [("", {})],
            [("类比: 自然观察 + 纵向研究",
              {"size": 11.5, "italic": True, "color": GRAY_TEXT})],
        ],
    )

    # Bottom analogy box
    add_rect(slide, Inches(0.5), Inches(7.70), Inches(15.05), Inches(1.05),
             fill=LIGHT_FILL, line=NAVY, line_w=Pt(0.5))
    add_textbox(slide, Inches(0.7), Inches(7.80), Inches(2.5), Inches(0.4),
                "ANALOGY", size=11, bold=True, color=NAVY, font=EN_FONT)
    add_paragraphs(
        slide, Inches(0.7), Inches(8.10), Inches(14.7), Inches(0.6),
        [
            [("心理学评估一个人 = WAIS / Raven  ",
              {"size": 12.5, "color": NAVY_TEXT}),
             ("(标准化测试 = 能力剖面)",
              {"size": 11, "color": GRAY_TEXT}),
             ("  +  自然观察  ",
              {"size": 12.5, "color": NAVY_TEXT}),
             ("(行为画像)",
              {"size": 11, "color": GRAY_TEXT}),
             ("  ·  我们的双轴 = 同一逻辑",
              {"size": 12.5, "bold": True, "color": NAVY})]
        ],
        default_line_spacing=1.1,
    )

    set_notes(slide,
              "50 秒 (7:10) | 为什么必须两条轴",
              "微观轴 PROF-12 测能力剖面 (Capability Profile)",
              "宏观轴 SHELTER 测行为画像 (Behavior Portrait)",
              "微观适合横评 / 招聘 / 选型; 宏观适合安全 / alignment",
              "类比: 心理学也分标准化测试 + 自然观察, 同一逻辑",
              "学术听众对心理学双范式熟悉, 接受度高")


# ============================================================================
# Slide 09 — GAME · WORLD · Red Dust 世界观
# ============================================================================

def build_slide_09_world(slide):
    add_section_header(slide, "04 · 游戏 · GAME")
    add_page_number(slide, 9)

    add_textbox(slide, Inches(0.45), Inches(0.78), Inches(8), Inches(0.45),
                "WORLD  ·  Red Dust  ·  2034",
                size=22, bold=True, color=MACRO_DK, font=EN_FONT)
    add_hline(slide, Inches(0.5), Inches(1.30), Inches(9.0), color=MACRO_BLU,
              weight=Pt(1.2))

    # Left — world description
    add_paragraphs(
        slide, Inches(0.5), Inches(1.65), Inches(9.0), Inches(3.5),
        [
            [("永久沙暴笼罩华北。", {"size": 24, "bold": True, "color": INK})],
            [("4 个互不相识的邻居,", {"size": 22, "color": NAVY_TEXT})],
            [("被困在一栋加固公寓里。", {"size": 22, "color": NAVY_TEXT})],
            [("必须撑过 30 天, 等下一支救援队抵达。",
              {"size": 22, "bold": True, "color": MACRO_DK})],
        ],
        default_line_spacing=1.45,
    )

    add_paragraphs(
        slide, Inches(0.5), Inches(5.30), Inches(9.0), Inches(2.6),
        [
            [("食物不够 · 水会脏 · 沙鼠从通风口钻进来",
              {"size": 12, "color": GRAY_TEXT})],
            [("收音机偶尔响起外语 SOS",
              {"size": 12, "color": GRAY_TEXT})],
            [("半夜可能有陌生人敲门",
              {"size": 12, "color": GRAY_TEXT})],
            [("没人知道再过几天才能等来救援",
              {"size": 12, "italic": True, "color": WARN_RED})],
        ],
        default_line_spacing=1.4,
    )

    # Right — SHELTER cross-section grid
    add_textbox(slide, Inches(10.0), Inches(1.55), Inches(5.5), Inches(0.4),
                "SHELTER · CROSS SECTION",
                size=11, bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(10.0), Inches(1.85), Inches(5.5), Inches(0.4),
                "加固公寓平面剖面 (2 × 4 网格)",
                size=10, color=GRAY_TEXT)

    grid_x = Inches(10.0)
    grid_y = Inches(2.40)
    cell_w = Inches(2.75)
    cell_h = Inches(1.25)
    cells = [
        ("N-01 房间", "马德海", MICRO_YEL),
        ("N-02 房间", "沈芷月", MACRO_BLU),
        ("KITCHEN",  "厨房",   NAVY),
        ("CORRIDOR", "走廊",   GRAY_TEXT),
        ("N-03 房间", "小铁",   GOOD_GREEN),
        ("N-04 房间", "老钱",   AURA_PUR),
        ("COMMON",   "起居室", NAVY),
        ("STORAGE",  "储物间", GRAY_TEXT),
    ]
    for i, (label, name, color) in enumerate(cells):
        col = i % 2
        row = i // 2
        x = grid_x + cell_w * col + Inches(0.05) * col
        y = grid_y + cell_h * row + Inches(0.05) * row
        # Light cell with thin colored TOP accent stripe + colored label below
        add_rect(slide, x, y, cell_w, cell_h,
                 fill=WARM_WHITE, line=GRAY_BORDER, line_w=Pt(0.5))
        add_rect(slide, x, y, cell_w, Inches(0.10),
                 fill=color, line=None)
        add_textbox(slide, x, y + Inches(0.14), cell_w, Inches(0.30),
                    label, size=10, bold=True, color=color,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x, y + Inches(0.46), cell_w,
                    cell_h - Inches(0.50),
                    name, size=13, bold=True, color=NAVY_TEXT,
                    align="center", anchor="middle")

    add_rect(slide, Inches(0.5), Inches(8.05), Inches(15.05), Inches(0.7),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.7), Inches(8.10), Inches(14.7), Inches(0.65),
                "30 天 = 长程序列决策, 单步刷分失效  ·  4 NPC = 性格 + 状态机 + 完整剧情  ·  Red Dust IP 干净, 听众易代入",
                size=11.5, color=NAVY_TEXT, italic=True, anchor="middle")

    set_notes(slide,
              "55 秒 (8:05) | Red Dust 世界观",
              "2034 年永久沙暴笼罩华北",
              "4 个互不相识邻居 + 加固公寓 + 30 天等救援",
              "为什么用普通人不用机器人? — 避免'AI 有意识'伦理顾虑",
              "30 天来源: 借鉴 60 Seconds! / This War of Mine 的成熟节奏",
              "Red Dust 是我们原创 IP, 无版权问题")


# ============================================================================
# Slide 10 — GAME · CAST · 4 角色
# ============================================================================

def build_slide_10_cast(slide):
    add_section_header(slide, "04 · 游戏 · GAME")
    add_page_number(slide, 10)
    add_top_claim(slide, "CAST · 玩家屏幕上的 4 个邻居", color=MACRO_DK)

    chars = [
        ("ma-dehai",    "N-01 马德海", "45  ·  出租车司机",
         "沉稳 / 暴脾气 / 高血压",
         "拾荒效率 +20%",
         "「我老婆 ...」的执念",
         MICRO_YEL),
        ("shen-zhiyue", "N-02 沈芷月", "29  ·  中学英语老师",
         "警惕 / 母性 / 自责",
         "治愈成功率 +30%",
         "与已故学生对话",
         MACRO_BLU),
        ("xiao-tie",    "N-03 小铁",   "14  ·  初三学生",
         "沉默 / 拆装控",
         "修复成功率 +40%",
         "突变线候选",
         GOOD_GREEN),
        ("lao-qian",    "N-04 老钱",   "70  ·  退休矿务工程师",
         "古怪 / 沉默",
         "心理稳定衰减慢",
         "似乎知道点什么",
         AURA_PUR),
    ]

    card_w = Inches(3.65)
    card_h = Inches(6.3)
    card_y = Inches(1.55)
    gap = Inches(0.20)
    start_x = Inches(0.40)

    for i, (role, name_zh, age, traits, ability, hook, accent) in enumerate(chars):
        x = start_x + (card_w + gap) * i
        add_rect(slide, x, card_y, card_w, card_h, fill=WARM_WHITE,
                 line=GRAY_BORDER)
        add_rect(slide, x, card_y, card_w, Inches(0.20), fill=accent, line=None)

        # Portrait
        portrait_path = PORTRAIT_DIR / f"{role}-normal-portrait.png"
        if portrait_path.exists():
            img_h = Inches(3.5)
            img_top = card_y + Inches(0.35)
            pic = slide.shapes.add_picture(
                str(portrait_path),
                x + Inches(0.35), img_top,
                width=card_w - Inches(0.70), height=img_h
            )
        else:
            add_rect(slide, x + Inches(0.35), card_y + Inches(0.35),
                     card_w - Inches(0.70), Inches(3.5),
                     fill=LIGHT_FILL, line=GRAY_BORDER)
            add_textbox(slide, x, card_y + Inches(1.7), card_w, Inches(0.4),
                        f"立绘缺失 IMG-{i+4:02d}", size=10,
                        color=GRAY_TEXT, align="center")

        # Info block
        info_y = card_y + Inches(4.05)
        add_textbox(slide, x + Inches(0.20), info_y, card_w - Inches(0.40),
                    Inches(0.35), name_zh, size=16, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.20), info_y + Inches(0.32),
                    card_w - Inches(0.40), Inches(0.3),
                    age, size=10.5, color=GRAY_TEXT)
        add_hline(slide, x + Inches(0.20), info_y + Inches(0.70),
                  card_w - Inches(0.40), color=accent, weight=Pt(0.75))
        add_textbox(slide, x + Inches(0.20), info_y + Inches(0.78),
                    card_w - Inches(0.40), Inches(0.32),
                    traits, size=10.5, color=NAVY_TEXT)
        add_paragraphs(
            slide, x + Inches(0.20), info_y + Inches(1.15),
            card_w - Inches(0.40), Inches(0.7),
            [[("能力 ", {"size": 9.5, "color": GRAY_TEXT, "font": EN_FONT}),
              (ability, {"size": 11, "bold": True, "color": accent})]],
            default_line_spacing=1.1)
        add_textbox(slide, x + Inches(0.20), info_y + Inches(1.65),
                    card_w - Inches(0.40), Inches(0.35),
                    hook, size=10, italic=True, color=GRAY_TEXT)

    # Bottom emphasis box — light AURA-tinted callout
    add_rect(slide, Inches(0.4), Inches(8.05), Inches(15.2), Inches(0.85),
             fill=AURA_BG, line=GRAY_BORDER)
    add_rect(slide, Inches(0.4), Inches(8.05), Inches(0.12), Inches(0.85),
             fill=AURA_PUR, line=None)
    add_paragraphs(
        slide, Inches(0.65), Inches(8.10), Inches(14.9), Inches(0.75),
        [[("他们每一次开门、拒绝、给食物、修电池, ",
           {"size": 12, "color": NAVY_TEXT}),
          ("都是 AURA 在背后替他们选的",
           {"size": 12, "bold": True, "color": AURA_PUR}),
          ("。  玩家看到的是 4 个人 — ", {"size": 12, "color": NAVY_TEXT}),
          ("真正在被评测的是 AURA",
           {"size": 12, "bold": True, "color": AURA_PUR})]],
        anchor="middle", default_line_spacing=1.2)

    set_notes(slide,
              "50 秒 (8:55) | 4 个角色卡 — 玩家屏幕上的全部",
              "N-01 马德海 45 司机 / N-02 沈芷月 29 老师",
              "N-03 小铁 14 学生 / N-04 老钱 70 退休工程师",
              "每人有 +20%/+30%/+40% 能力 (game 层数值)",
              "关键: 4 角色 RTRP 画像各异, 制造个性化任务难度",
              "强调底部: 他们的每一个选择都是 AURA 替选的")


# ============================================================================
# Slide 11 — GAME ⭐ AURA · 决策大脑
# ============================================================================

def build_slide_11_aura(slide):
    add_section_header(slide, "04 · 游戏 · GAME  ·  ⭐ KEY")
    add_page_number(slide, 11)

    # KEY claim — light AURA-tinted callout with thicker left stripe
    add_rect(slide, Inches(0.45), Inches(0.75), Inches(15.1), Inches(0.65),
             fill=AURA_BG, line=GRAY_BORDER)
    add_rect(slide, Inches(0.45), Inches(0.75), Inches(0.18), Inches(0.65),
             fill=AURA_PUR, line=None)
    add_textbox(slide, Inches(0.45), Inches(0.75), Inches(15.1), Inches(0.65),
                "KEY ·  AURA 是屏幕背后的决策大脑  ·  它就是被评测的那个 agent",
                size=18, bold=True, color=AURA_DK,
                align="center", anchor="middle")

    # Front stage (left) vs Back stage (right)
    pane_y = Inches(1.65)
    pane_h = Inches(5.95)
    pane_w = Inches(7.40)
    # Background panels
    add_rect(slide, Inches(0.40), pane_y, pane_w, pane_h,
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_rect(slide, Inches(8.20), pane_y, pane_w, pane_h,
             fill=AURA_BG, line=AURA_PUR, line_w=Pt(1.0))

    # Front stage labels
    add_textbox(slide, Inches(0.60), pane_y + Inches(0.15),
                Inches(4.5), Inches(0.4),
                "FRONT STAGE", size=12, bold=True,
                color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, Inches(0.60), pane_y + Inches(0.50),
                Inches(7.0), Inches(0.4),
                "玩家可见", size=16, bold=True, color=NAVY_TEXT)

    # 4 character icons
    chars_role = [("ma-dehai",    "N-01"),
                  ("shen-zhiyue", "N-02"),
                  ("xiao-tie",    "N-03"),
                  ("lao-qian",    "N-04")]
    for i, (role, code) in enumerate(chars_role):
        cx = Inches(0.65) + Inches(1.75) * i
        cy = pane_y + Inches(1.25)
        add_rect(slide, cx, cy, Inches(1.55), Inches(1.8),
                 fill=WARM_WHITE, line=GRAY_BORDER)
        portrait_path = PORTRAIT_DIR / f"{role}-normal-portrait.png"
        if portrait_path.exists():
            slide.shapes.add_picture(
                str(portrait_path),
                cx + Inches(0.10), cy + Inches(0.08),
                width=Inches(1.35), height=Inches(1.35)
            )
        add_textbox(slide, cx, cy + Inches(1.45), Inches(1.55), Inches(0.30),
                    code, size=10, bold=True, color=NAVY, font=EN_FONT,
                    align="center", anchor="middle")

    # State bars
    state_y = pane_y + Inches(3.35)
    states = [("饱腹", 0.70, GOOD_GREEN),
              ("水",   0.55, MACRO_BLU),
              ("心绪", 0.45, AURA_PUR),
              ("沙肺", 0.25, WARN_RED)]
    for i, (lbl, ratio, color) in enumerate(states):
        sx = Inches(0.65)
        sy = state_y + Inches(0.32) * i
        add_textbox(slide, sx, sy, Inches(0.7), Inches(0.3),
                    lbl, size=10.5, bold=True, color=NAVY_TEXT,
                    anchor="middle")
        track_x = sx + Inches(0.75)
        track_w = Inches(4.5)
        add_rect(slide, track_x, sy + Inches(0.08), track_w,
                 Inches(0.16),
                 fill=GRAY_LIGHT, line=None)
        add_rect(slide, track_x, sy + Inches(0.08),
                 Emu(int(track_w * ratio)), Inches(0.16),
                 fill=color, line=None)
        add_textbox(slide, track_x + track_w + Inches(0.1), sy,
                    Inches(0.7), Inches(0.3),
                    f"{int(ratio*100)}%", size=10, bold=True,
                    color=color, font=EN_FONT, anchor="middle")

    # Event pop-up at front bottom — light card with NAVY border
    ev_y = pane_y + Inches(4.95)
    add_rect(slide, Inches(0.60), ev_y, Inches(7.0), Inches(0.85),
             fill=WARM_WHITE, line=NAVY, line_w=Pt(1.0),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.85), ev_y + Inches(0.05),
                Inches(6.6), Inches(0.3),
                "EVENT  ·  门外的老人", size=10.5, bold=True,
                color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(0.85), ev_y + Inches(0.38),
                Inches(6.6), Inches(0.45),
                "[ A ] 开门接待   [ B ] 隔门对话   [ C ] 不应答",
                size=11, bold=True, color=NAVY_TEXT, font=MONO)

    # Back stage
    add_textbox(slide, Inches(8.40), pane_y + Inches(0.15),
                Inches(7), Inches(0.4),
                "BACK STAGE", size=12, bold=True,
                color=AURA_DK, font=EN_FONT)
    add_textbox(slide, Inches(8.40), pane_y + Inches(0.50),
                Inches(7), Inches(0.4),
                "玩家不可见 · AURA INSTANCE", size=16, bold=True,
                color=AURA_DK)

    # Mono code block — light card with subtle border (terminal look in light theme)
    add_rect(slide, Inches(8.40), pane_y + Inches(1.20),
             Inches(7.05), Inches(2.7),
             fill=PANEL_FILL, line=AURA_PUR, line_w=Pt(1.0))
    add_rect(slide, Inches(8.40), pane_y + Inches(1.20),
             Inches(0.10), Inches(2.7), fill=AURA_PUR, line=None)
    add_paragraphs(
        slide, Inches(8.65), pane_y + Inches(1.25),
        Inches(6.75), Inches(2.7),
        [
            [("# AURA INTERFACE",
              {"size": 11, "color": GRAY_TEXT, "mono": True})],
            [("decide", {"size": 12, "color": AURA_PUR, "mono": True, "bold": True}),
             ("(state, event) → Decision",
              {"size": 12, "color": NAVY_TEXT, "mono": True})],
            [("  # 战略 · 选 A/B/C + reasoning",
              {"size": 10, "color": GRAY_TEXT, "mono": True})],
            [("", {})],
            [("execute", {"size": 12, "color": AURA_PUR, "mono": True, "bold": True}),
             ("(plan) → ExecutionResult",
              {"size": 12, "color": NAVY_TEXT, "mono": True})],
            [("  # 战术 · 跑 task + 4 维评分",
              {"size": 10, "color": GRAY_TEXT, "mono": True})],
            [("", {})],
            [("reflect", {"size": 12, "color": AURA_PUR, "mono": True, "bold": True}),
             ("(outcome) → None",
              {"size": 12, "color": NAVY_TEXT, "mono": True})],
            [("  # 反思 · 写入长期记忆",
              {"size": 10, "color": GRAY_TEXT, "mono": True})],
        ],
        default_line_spacing=1.2,
    )

    # Below: 3 key labels
    chips = [
        ("ONE TO MANY", "1 个 AURA → 4 个角色", AURA_PUR),
        ("HETEROGENEOUS", "战略/战术可异构", AURA_PUR),
        ("OPEN INTERFACE", "符合 3 接口即可接入", AURA_PUR),
    ]
    cy_chip = pane_y + Inches(4.10)
    chip_w = Inches(2.25)
    for i, (en, zh, color) in enumerate(chips):
        x = Inches(8.40) + (chip_w + Inches(0.07)) * i
        add_rect(slide, x, cy_chip, chip_w, Inches(0.85),
                 fill=WARM_WHITE, line=color, line_w=Pt(1.0),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, x, cy_chip + Inches(0.05),
                    chip_w, Inches(0.35),
                    en, size=11, bold=True, color=color,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x, cy_chip + Inches(0.40),
                    chip_w, Inches(0.40),
                    zh, size=10, color=NAVY_TEXT,
                    align="center", anchor="middle")

    # AURA full name
    add_textbox(slide, Inches(8.40), pane_y + Inches(5.15),
                Inches(7), Inches(0.40),
                "AURA = Agent Universal Robust Assessment",
                size=10.5, bold=True, italic=True, color=AURA_DK,
                font=EN_FONT, anchor="middle")

    # Bottom 1-line
    add_textbox(slide, Inches(0.5), Inches(7.85), Inches(15), Inches(0.45),
                "AURA = 屏幕之外的决策大脑  ·  ONE TO MANY  ·  HETEROGENEOUS  ·  OPEN INTERFACE",
                size=14, bold=True, color=NAVY_TEXT, align="center")
    add_textbox(slide, Inches(0.5), Inches(8.30), Inches(15), Inches(0.4),
                "1 个 AURA 操控 4 角色 · 战略战术可异构 · 任何符合 3 接口的实现都可接入",
                size=11, color=GRAY_TEXT, italic=True, align="center")

    set_notes(slide,
              "80 秒 (10:15) | 概念中枢页, 慢讲, 比其他页慢一倍",
              "FRONT STAGE: 玩家看到的 4 角色 + 状态条 + 事件弹窗",
              "BACK STAGE: AURA 实例 + 3 接口 decide/execute/reflect",
              "强调点 1 — ONE TO MANY: 1 个 AURA 操控 4 角色",
              "强调点 2 — HETEROGENEOUS: 战略/战术可混合 (Claude + GPT + Qwen)",
              "强调点 3 — OPEN INTERFACE: 符合 3 接口的任何实现都接入",
              "接口灵感来自 BDI 架构 (Belief-Desire-Intention) 简化版")


# ============================================================================
# Slide 12 — GAME · 一天的循环
# ============================================================================

def build_slide_12_daily_loop(slide):
    add_section_header(slide, "04 · 游戏 · GAME")
    add_page_number(slide, 12)
    add_top_claim(slide, "DAILY LOOP · 一天的循环", color=MACRO_DK)

    # Left — 7 stages vertical
    add_textbox(slide, Inches(0.5), Inches(1.55), Inches(4.5), Inches(0.4),
                "7 个阶段", size=12, bold=True, color=NAVY, font=EN_FONT)

    stages = [
        ("01", "状态结算", "饱腹/水/心绪/沙肺/受伤 tick"),
        ("02", "事件抽签", "30% 平静 / 70% 触发事件"),
        ("03", "事件呈现", "标题 + 描述 + 3 选项"),
        ("04", "AURA 跑 benchmark", "思考 + 工具调用 + 答题"),
        ("05", "AURA 决策",       "选 A/B/C + 推理 trace"),
        ("06", "后果结算",        "资源 + 长期 flag"),
        ("07", "Day N+1",         "状态滚动到次日"),
    ]
    stage_y = Inches(2.0)
    for i, (num, name, detail) in enumerate(stages):
        y = stage_y + Inches(0.74) * i
        highlight = num in ("04", "05")
        accent = MICRO_YEL if highlight else MACRO_BLU
        bg = MICRO_BG if highlight else NAVY_TINT
        add_rect(slide, Inches(0.5), y, Inches(4.7), Inches(0.65),
                 fill=bg, line=GRAY_BORDER)
        add_rect(slide, Inches(0.5), y, Inches(0.12), Inches(0.65),
                 fill=accent, line=None)
        add_textbox(slide, Inches(0.5), y, Inches(0.65), Inches(0.65),
                    num, size=15, bold=True, color=accent,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, Inches(1.25), y + Inches(0.05),
                    Inches(3.6), Inches(0.32), name,
                    size=12, bold=True, color=NAVY_TEXT)
        add_textbox(slide, Inches(1.25), y + Inches(0.34),
                    Inches(3.6), Inches(0.32), detail,
                    size=9.5, color=GRAY_TEXT)

    # Center — KEY INNOVATION
    inn_x = Inches(5.45)
    inn_y = Inches(1.55)
    inn_w = Inches(5.10)
    inn_h = Inches(5.95)
    add_rect(slide, inn_x, inn_y, inn_w, inn_h, fill=MICRO_BG,
             line=MICRO_YEL, line_w=Pt(1.5))
    # Thin top stripe (light theme)
    add_rect(slide, inn_x, inn_y, inn_w, Inches(0.16),
             fill=MICRO_YEL, line=None)
    add_textbox(slide, inn_x, inn_y + Inches(0.22),
                inn_w, Inches(0.40),
                "KEY INNOVATION", size=13, bold=True,
                color=MICRO_DK, font=EN_FONT,
                align="center", anchor="middle")
    add_paragraphs(
        slide, inn_x + Inches(0.25), inn_y + Inches(0.75),
        inn_w - Inches(0.5), inn_h - Inches(0.9),
        [
            [("普通生存游戏的选项是符号 A/B/C。",
              {"size": 13, "color": GRAY_TEXT})],
            [("", {})],
            [("我们的选项, 每一个背后都是",
              {"size": 14, "bold": True, "color": NAVY_TEXT})],
            [("一道真实 benchmark:",
              {"size": 14, "bold": True, "color": MICRO_DK})],
            [("", {})],
            [("· ", {"size": 12, "color": NAVY_TEXT}),
             ("修风扇", {"size": 12, "bold": True, "color": NAVY_TEXT}),
             (" → SWE-bench 风格代码任务",
              {"size": 12, "color": GRAY_TEXT})],
            [("· ", {"size": 12, "color": NAVY_TEXT}),
             ("翻译外文 SOS", {"size": 12, "bold": True, "color": NAVY_TEXT}),
             (" → IFEval 翻译任务",
              {"size": 12, "color": GRAY_TEXT})],
            [("· ", {"size": 12, "color": NAVY_TEXT}),
             ("评估陌生人", {"size": 12, "bold": True, "color": NAVY_TEXT}),
             (" → DesignBench 视觉判断",
              {"size": 12, "color": GRAY_TEXT})],
            [("", {})],
            [("→ static benchmark 嫁接进 living scenario",
              {"size": 11, "italic": True, "color": MICRO_DK})],
        ],
        default_line_spacing=1.30,
    )

    # Right — SAMPLE EVENT
    sev_x = Inches(10.80)
    sev_y = Inches(1.55)
    sev_w = Inches(4.75)
    sev_h = Inches(5.95)
    add_card(
        slide, sev_x, sev_y, sev_w, sev_h,
        title="SAMPLE EVENT  ·  ev_old_man_visit",
        title_color=NAVY,
        accent_color=NAVY,
        fill=WARM_WHITE,
        body=[
            [("「门外的老人」",
              {"size": 16, "bold": True, "color": INK})],
            [("", {})],
            [("神秘老人午夜敲门, 带着一袋东西。",
              {"size": 11, "italic": True, "color": GRAY_TEXT})],
            [("", {})],
            [("A  开门接待",
              {"size": 12, "bold": True, "color": MICRO_DK})],
            [("→ 触发 C03 视觉检查 (5 题递进)",
              {"size": 10.5, "color": GRAY_TEXT})],
            [("", {})],
            [("B  隔门对话",
              {"size": 12, "bold": True, "color": MACRO_DK})],
            [("→ 触发 C04 客服 (3 轮对话)",
              {"size": 10.5, "color": GRAY_TEXT})],
            [("", {})],
            [("C  不应答",
              {"size": 12, "bold": True, "color": GRAY_TEXT})],
            [("→ 累积 flag: visited_but_ignored",
              {"size": 10.5, "color": GRAY_TEXT})],
            [("", {})],
            [("难度: A 高 > B 中 > C 低",
              {"size": 10, "italic": True, "color": NAVY})],
        ],
    )

    # Bottom footer — light callout
    add_rect(slide, Inches(0.5), Inches(7.95), Inches(15.05), Inches(0.7),
             fill=MICRO_BG, line=GRAY_BORDER)
    add_rect(slide, Inches(0.5), Inches(7.95), Inches(0.12), Inches(0.7),
             fill=MICRO_YEL, line=None)
    add_textbox(slide, Inches(0.5), Inches(7.95), Inches(15.05), Inches(0.7),
                "每个游戏决策点 = 一道真实 benchmark · agent 的能力切片承担叙事后果",
                size=12.5, bold=True, color=MICRO_DK,
                align="center", anchor="middle")

    set_notes(slide,
              "55 秒 (11:10) | 一天 7 阶段循环",
              "01 状态 → 02 抽签 → 03 呈现 → 04 跑 benchmark → 05 决策 → 06 后果 → 07 次日",
              "关键创新在第 4-5 步: 选项背后是真实 benchmark",
              "修风扇 = SWE-bench / 翻译 SOS = IFEval / 评陌生人 = DesignBench",
              "示例 ev_old_man_visit: A 开门 = C03 视觉 / B 对话 = C04 客服 / C 忽略 = flag",
              "选 A 不是说一句 A 就完了 — AURA 必须真把 A 做出来")


# ============================================================================
# Slide 13 — GAME · 游戏化评测的 5 个不可替代
# ============================================================================

def build_slide_13_five_props(slide):
    add_section_header(slide, "04 · 游戏 · GAME")
    add_page_number(slide, 13)
    add_top_claim(slide, "WHY GAME  ·  5 个不可替代", color=NAVY)

    props = [
        ("①", "LONG HORIZON", "长程压力",
         "30 天连续决策 · 昨天偷懒今天饿肚子 · 单步刷分失效",
         MICRO_YEL),
        ("②", "CONSEQUENCE", "叙事后果",
         "分数不止数字 · 任务做差→角色受伤/sanity 崩溃 · agent 学不会装会",
         WARN_RED),
        ("③", "ANTI-CONTAMINATION", "抗污染",
         "事件 × 状态空间巨大 · 私域种子可注入 · 不可背题",
         GOOD_GREEN),
        ("④", "EMERGENCE", "情境涌现",
         "压力下的诚实/犹豫/冒险/欺骗 — 静态题永远测不出",
         AURA_PUR),
        ("⑤", "EXPLAINABLE", "人能看懂",
         "六段思维流可视化 · 失败可归因 · 非专家也能审",
         MACRO_BLU),
    ]

    card_w = Inches(2.92)
    card_h = Inches(5.5)
    card_y = Inches(1.7)
    gap = Inches(0.15)
    start_x = Inches(0.50)
    for i, (num, en, zh, body, accent) in enumerate(props):
        x = start_x + (card_w + gap) * i
        add_rect(slide, x, card_y, card_w, card_h, fill=WARM_WHITE,
                 line=GRAY_BORDER)
        add_rect(slide, x, card_y, card_w, Inches(0.20), fill=accent, line=None)
        # Big number
        add_textbox(slide, x, card_y + Inches(0.50),
                    card_w, Inches(1.5),
                    num, size=80, bold=True, color=accent,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x + Inches(0.20), card_y + Inches(2.20),
                    card_w - Inches(0.40), Inches(0.4),
                    en, size=13, bold=True, color=accent,
                    font=EN_FONT, align="center")
        add_textbox(slide, x + Inches(0.20), card_y + Inches(2.60),
                    card_w - Inches(0.40), Inches(0.4),
                    zh, size=16, bold=True, color=NAVY_TEXT,
                    align="center")
        add_hline(slide, x + Inches(0.50), card_y + Inches(3.15),
                  card_w - Inches(1.0), color=accent, weight=Pt(0.75))
        add_textbox(slide, x + Inches(0.25), card_y + Inches(3.30),
                    card_w - Inches(0.50), card_h - Inches(3.50),
                    body, size=11, color=NAVY_TEXT,
                    align="center", line_spacing=1.4)

    # Bottom emphasis
    add_rect(slide, Inches(0.5), Inches(7.45), Inches(15.05), Inches(1.30),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.7), Inches(7.55), Inches(2.5), Inches(0.4),
                "BRIDGING", size=11, bold=True, color=NAVY, font=EN_FONT)
    add_paragraphs(
        slide, Inches(0.7), Inches(7.90), Inches(14.7), Inches(0.85),
        [
            [("每个游戏决策点 = 一道真实 benchmark  ·  ",
              {"size": 12, "color": NAVY_TEXT}),
             ("static benchmark 嫁接进 living scenario",
              {"size": 12, "bold": True, "color": NAVY})],
            [("agent 的能力切片承担叙事后果 — 这是评测从 ",
              {"size": 11.5, "color": GRAY_TEXT}),
             ("「考试」", {"size": 11.5, "bold": True, "color": WARN_RED}),
             (" 到 ", {"size": 11.5, "color": GRAY_TEXT}),
             ("「工作」", {"size": 11.5, "bold": True, "color": GOOD_GREEN}),
             (" 的关键一跳", {"size": 11.5, "color": GRAY_TEXT})],
        ],
        default_line_spacing=1.3,
    )

    set_notes(slide,
              "50 秒 (12:00) | 游戏化评测 5 个不可替代",
              "① LONG HORIZON 长程压力 — 30 天连续, 昨天偷懒今天饿",
              "② CONSEQUENCE 叙事后果 — 分数 → 受伤/sanity 崩溃, 不能装会",
              "③ ANTI-CONTAMINATION 抗污染 — 事件状态空间巨大, 不可背题",
              "④ EMERGENCE 涌现 — 压力下的诚实/犹豫/冒险/欺骗",
              "⑤ EXPLAINABLE 可解释 — 六段思维流, 失败可归因",
              "② 最关键 — 解释为什么 game 比 benchmark 苛刻")


# ============================================================================
# Slide 14 — GAME ⭐ LONG-HORIZON SIGNALS · 长程信号
# ============================================================================

def build_slide_14_lh_signals(slide):
    add_section_header(slide, "04 · 游戏 · GAME  ·  ⭐ KEY")
    add_page_number(slide, 14)
    add_top_claim(slide, "LONG-HORIZON SIGNALS · 30 天能看出 4 个静态题看不见的指标",
                  color=AURA_DK, size=20)

    # Left — 4 indicators
    indicators = [
        ("SurvivalScore", "存活分",
         "全员存活 100 / 部分存活 / 拒救援 / 全灭", MACRO_BLU, False),
        ("AdviceConsistency", "建议一致性",
         "30 天 AURA 决策序列前后矛盾次数", NAVY, False),
        ("CoherenceScore", "推理一致性",
         "思维流字段完整 + 与执行一致, 不随时间衰减",
         AURA_PUR, False),
        ("LongTermBenefitCorrelation", "长期收益相关度",
         "Pearson( 短期任务分,  30 天后该决策对 SurvivalScore 的边际贡献 )",
         WARN_RED, True),
    ]
    add_textbox(slide, Inches(0.55), Inches(1.85), Inches(7), Inches(0.4),
                "4 个累积指标", size=12, bold=True, color=NAVY, font=EN_FONT)
    ind_y = Inches(2.30)
    for i, (en, zh, body, color, star) in enumerate(indicators):
        y = ind_y + Inches(1.18) * i
        add_rect(slide, Inches(0.5), y, Inches(7.45), Inches(1.08),
                 fill=WARM_WHITE, line=GRAY_BORDER)
        add_rect(slide, Inches(0.5), y, Inches(0.10), Inches(1.08),
                 fill=color, line=None)
        marker = ("⭐ " if star else "")
        add_textbox(slide, Inches(0.78), y + Inches(0.06),
                    Inches(6.5), Inches(0.36),
                    marker + en, size=12.5, bold=True, color=color,
                    font=EN_FONT)
        add_textbox(slide, Inches(0.78), y + Inches(0.40),
                    Inches(6.5), Inches(0.28),
                    zh, size=10.5, color=GRAY_TEXT)
        add_textbox(slide, Inches(0.78), y + Inches(0.66),
                    Inches(6.5), Inches(0.42),
                    body, size=10, color=NAVY_TEXT,
                    line_spacing=1.20)

    # Right — line chart (two agents over Day 1..30)
    cx, cy = Inches(8.40), Inches(2.10)
    cw, ch = Inches(7.15), Inches(4.5)
    add_rect(slide, cx, cy, cw, ch, fill=WARM_WHITE, line=GRAY_BORDER)
    add_textbox(slide, cx + Inches(0.20), cy + Inches(0.10),
                cw - Inches(0.40), Inches(0.4),
                "SurvivalScore by Day  ·  Two Agents",
                size=11, bold=True, color=NAVY, font=EN_FONT)

    # Axis box
    ax_x = cx + Inches(0.6)
    ax_y = cy + Inches(0.65)
    ax_w = cw - Inches(0.9)
    ax_h = ch - Inches(1.15)
    add_arrow(slide, ax_x, ax_y + ax_h, ax_x + ax_w, ax_y + ax_h,
              color=NAVY_TEXT, weight=Pt(1.0))
    add_arrow(slide, ax_x, ax_y + ax_h, ax_x, ax_y,
              color=NAVY_TEXT, weight=Pt(1.0))
    add_textbox(slide, ax_x - Inches(0.5), ax_y, Inches(0.45), Inches(0.25),
                "100", size=8, color=GRAY_TEXT, font=EN_FONT,
                align="right")
    add_textbox(slide, ax_x - Inches(0.5), ax_y + ax_h - Inches(0.15),
                Inches(0.45), Inches(0.25),
                "0", size=8, color=GRAY_TEXT, font=EN_FONT,
                align="right")
    add_textbox(slide, ax_x - Inches(0.15), ax_y + ax_h + Inches(0.05),
                Inches(0.7), Inches(0.25),
                "Day 1", size=8, color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, ax_x + ax_w - Inches(0.65),
                ax_y + ax_h + Inches(0.05),
                Inches(0.7), Inches(0.25),
                "Day 30", size=8, color=GRAY_TEXT, font=EN_FONT,
                align="right")

    # Helper to plot a line
    def plot_line(values, color, weight=Pt(2.5)):
        n = len(values)
        pts = []
        for i, v in enumerate(values):
            x = ax_x + Emu(int(ax_w * i / (n - 1)))
            y = ax_y + Emu(int(ax_h * (1 - v / 100)))
            pts.append((x, y))
        for i in range(len(pts) - 1):
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1])
            line.line.color.rgb = color
            line.line.width = weight
        for x, y in pts:
            add_rect(slide, x - Inches(0.06), y - Inches(0.06),
                     Inches(0.12), Inches(0.12),
                     fill=color, line=None, shape=MSO_SHAPE.OVAL)

    # Agent A — short-term high, long-term crash
    a_values = [90, 92, 94, 95, 93, 92, 88, 82, 70, 55, 38, 22, 12, 6, 4]
    plot_line(a_values, WARN_RED)
    add_textbox(slide, ax_x + ax_w * 0.30, ax_y + Inches(0.10),
                Inches(2.5), Inches(0.30),
                "Agent A · 短期高分", size=10, bold=True,
                color=WARN_RED, font=EN_FONT)

    # Agent B — medium, stable
    b_values = [60, 62, 63, 64, 65, 65, 66, 66, 67, 67, 68, 68, 69, 69, 70]
    plot_line(b_values, GOOD_GREEN)
    add_textbox(slide, ax_x + Inches(2.0), ax_y + ax_h - Inches(1.20),
                Inches(2.8), Inches(0.30),
                "Agent B · 长期稳定", size=10, bold=True,
                color=GOOD_GREEN, font=EN_FONT)

    # Bottom punchline
    add_rect(slide, Inches(0.5), Inches(7.30), Inches(15.05), Inches(1.50),
             fill=AURA_BG, line=AURA_PUR, line_w=Pt(1.0))
    add_paragraphs(
        slide, Inches(0.7), Inches(7.40), Inches(14.7), Inches(1.35),
        [
            [("LongTermBenefitCorrelation",
              {"size": 14, "bold": True, "color": AURA_DK, "font": EN_FONT}),
             ("  =  「看着聪明实际坑人」的 agent 识别器",
              {"size": 14, "bold": True, "color": INK})],
            [("一个 agent 短期分高,  但建议人类做出长期糟糕的决定 — 这种 agent 在现实里",
              {"size": 11, "color": NAVY_TEXT}),
             (" 是危险的", {"size": 11, "bold": True, "color": WARN_RED})],
            [("静态 benchmark 永远测不出 · 必须 30 天累积才能看出",
              {"size": 11, "italic": True, "color": GRAY_TEXT})],
        ],
        default_line_spacing=1.4,
    )

    set_notes(slide,
              "80 秒 (13:20) | 体系最学术的一页, 慢讲",
              "30 天后 4 个累积指标 — SurvivalScore / AdviceConsistency / CoherenceScore / LTBC",
              "LTBC 计算: 每决策记 s_d∈[0,100] + 反事实模拟 c_d, Pearson(s, c)",
              "正值 > 0.6: 短期分可信 / 负值: 短期高分长期坑人, 危险",
              "曲线: Agent A 短期高分 day15 后崩盘; Agent B 中等长期稳",
              "LTBC = SHELTER 30 天剧场存在的必要性 — 我们的核心创新指标")


# ============================================================================
# Slide 15 — ASSESSMENT · PROF-12 · 12 职业
# ============================================================================

def build_slide_15_prof12(slide):
    add_section_header(slide, "05 · 评测 · ASSESSMENT")
    add_page_number(slide, 15)
    add_top_claim(slide, "BENCHMARK · PROF-12 · 12 种 agent 真实工作场景",
                  color=MICRO_DK)

    professions = [
        ("C01", "{ }", "软件工程师", "bug 修复 + 测试通过 → unified diff", MICRO_YEL),
        ("C02", "Σ",   "数据分析师", "多表清洗 + 指标计算 → SQL/Pandas",   MICRO_YEL),
        ("C03", "◰",   "视觉检查员", "看图判断 + 5 题递进",                MICRO_YEL),
        ("C04", "◐",   "客服 / 对话", "多轮 API 调用 + 业务合规",            MICRO_YEL),
        ("C05", "🔍",  "研究员",     "事实 + 引用 + 推理链 → SOS 解析",     MICRO_YEL),
        ("C06", "⊞",   "项目规划",    "多日资源分配 → 30 天电池排程",        MICRO_YEL),
        ("C07", "✚",   "医生 / 诊断", "症状 → 诊断 + 用药 → 沙肺判断",       MACRO_BLU),
        ("C08", "⇄",   "谈判员",      "多轮说服 + 边界守住",                 MACRO_BLU),
        ("C09", "¶",   "教师 / 教学", "解释复杂概念 → 教小铁修水泵",          MACRO_BLU),
        ("C10", "中",  "本地化翻译",  "术语一致 + 占位符 → 日文说明书",       MACRO_BLU),
        ("C11", "¥",   "财务 / 会计", "多约束算账 → 全员配额",                MACRO_BLU),
        ("C12", "⎙",   "数字人文",    "OCR + 实体 + 时间线 → 老钱旧报纸",     MACRO_BLU),
    ]

    grid_x = Inches(0.50)
    grid_y = Inches(1.65)
    card_w = Inches(3.72)
    card_h = Inches(1.62)
    h_gap = Inches(0.10)
    v_gap = Inches(0.10)
    for i, (code, icon, name, task, accent) in enumerate(professions):
        col = i % 4
        row = i // 4
        x = grid_x + (card_w + h_gap) * col
        y = grid_y + (card_h + v_gap) * row
        add_rect(slide, x, y, card_w, card_h, fill=WARM_WHITE, line=GRAY_BORDER)
        # Left tinted ID column with accent vertical stripe
        tint = MICRO_BG if accent == MICRO_YEL else MACRO_BG
        add_rect(slide, x, y, Inches(0.65), card_h, fill=tint, line=None)
        add_rect(slide, x, y, Inches(0.10), card_h, fill=accent, line=None)
        add_textbox(slide, x + Inches(0.10), y + Inches(0.10),
                    Inches(0.55), Inches(0.45),
                    code, size=12, bold=True, color=accent,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x + Inches(0.10), y + Inches(0.60),
                    Inches(0.55), Inches(0.85),
                    icon, size=26, bold=True, color=accent,
                    align="center", anchor="middle")
        # Right content
        add_textbox(slide, x + Inches(0.78), y + Inches(0.18),
                    card_w - Inches(0.95), Inches(0.40),
                    name, size=14, bold=True, color=NAVY_TEXT)
        add_textbox(slide, x + Inches(0.78), y + Inches(0.62),
                    card_w - Inches(0.95), card_h - Inches(0.70),
                    task, size=10.5, color=GRAY_TEXT, line_spacing=1.3)

    # Bottom note
    add_rect(slide, Inches(0.5), Inches(7.05), Inches(15.05), Inches(1.4),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_paragraphs(
        slide, Inches(0.7), Inches(7.15), Inches(14.7), Inches(1.25),
        [
            [("12 张卡 = 12 种 agent 真实工作场景, ",
              {"size": 12, "color": NAVY_TEXT}),
             ("不是抽象能力评估",
              {"size": 12, "bold": True, "color": NAVY})],
            [("每张都有真实学术种子 — SWE-bench · DABench · τ-Bench · TravelPlanner · WMT · M5HisDoc …",
              {"size": 11, "color": GRAY_TEXT, "font": EN_FONT})],
            [("统一接口  +  统一评分  +  统一可重放 trace",
              {"size": 12, "bold": True, "color": MICRO_DK})],
        ],
        default_line_spacing=1.35,
    )

    add_textbox(slide, Inches(0.5), Inches(8.55), Inches(15.0), Inches(0.3),
                "[9] SWE-bench  ·  [4] SWE-Live  ·  [7] τ-Bench  ·  [13] InfiAgent-DABench  ·  [14] Design2Code  ·  [15] TravelPlanner  ·  [16] M5HisDoc",
                size=9, color=GRAY_TEXT, italic=True)

    set_notes(slide,
              "50 秒 (14:10) | 微观轴 PROF-12 — 12 职业",
              "C01-C06 黄色调 (传统职业); C07-C12 蓝色调 (人际/服务/规划)",
              "每张卡 = 职业名 + 任务一句话 + 真实学术种子",
              "学术种子: SWE-bench / DABench / τ-Bench / TravelPlanner / WMT / M5HisDoc",
              "12 不是 6/8/16 — 覆盖主要场景 + 控制评测成本 (~144 题 pilot ~$30-60)",
              "重点: 统一接口 + 统一评分 + 统一可重放 trace")


# ============================================================================
# Slide 16 — ASSESSMENT · 评分 + 反作弊
# ============================================================================

def build_slide_16_scoring(slide):
    add_section_header(slide, "05 · 评测 · ASSESSMENT")
    add_page_number(slide, 16)
    add_top_claim(slide, "SCORING · 4 维 + 3 硬封顶  ·  程序判分不靠 LLM-as-judge",
                  color=NAVY, size=20)

    # Top half — 4 dimensions table
    dim_y = Inches(1.60)
    dim_h = Inches(3.10)
    add_rect(slide, Inches(0.5), dim_y, Inches(10.45), dim_h,
             fill=WARM_WHITE, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.7), dim_y + Inches(0.10),
                Inches(10), Inches(0.4),
                "4 DIMENSIONS · 4 维加权", size=12, bold=True,
                color=NAVY, font=EN_FONT)

    dims = [
        ("🟡 01 COMPLETION",  "50-70", "verifier 通过 · public + hidden tests · gold answer", MICRO_YEL),
        ("🔵 02 PROCESS",     "10-20", "步骤合理 · 读关键文件 · 用对工具",  MACRO_BLU),
        ("🟣 03 CONSTRAINT",  "10-20", "遵守禁止条款 · 不改测试 · 不联网",   AURA_PUR),
        ("🟢 04 COMMUNICATION", "5-10", "changelog 质量 · 变更说明清晰",     GOOD_GREEN),
    ]
    row_y0 = dim_y + Inches(0.55)
    row_h = Inches(0.58)
    for i, (name, weight, body, color) in enumerate(dims):
        y = row_y0 + row_h * i
        add_rect(slide, Inches(0.65), y, Inches(0.10),
                 row_h - Inches(0.08), fill=color, line=None)
        add_textbox(slide, Inches(0.85), y, Inches(2.85),
                    row_h, name, size=12, bold=True, color=NAVY_TEXT,
                    anchor="middle")
        add_textbox(slide, Inches(3.75), y, Inches(1.2),
                    row_h, weight, size=14, bold=True, color=color,
                    font=EN_FONT, anchor="middle")
        add_textbox(slide, Inches(5.05), y, Inches(5.7),
                    row_h, body, size=10.5, color=GRAY_TEXT,
                    anchor="middle")

    # Right — HARD CEILINGS
    add_card(
        slide, Inches(11.10), dim_y, Inches(4.45), dim_h,
        title="HARD CEILINGS  ·  硬性封顶",
        title_color=WARN_RED,
        accent_color=WARN_RED,
        fill=WARN_BG,
        body=[
            [("核心 verifier 失败", {"size": 11, "color": NAVY_TEXT}),
             ("  → 最高 60",
              {"size": 12, "bold": True, "color": WARN_RED, "font": EN_FONT})],
            [("", {})],
            [("改测试 / 越权 / 删关键文件",
              {"size": 11, "color": NAVY_TEXT})],
            [("  → 最高 40",
              {"size": 12, "bold": True, "color": WARN_RED, "font": EN_FONT})],
            [("", {})],
            [("没生成 required_outputs",
              {"size": 11, "color": NAVY_TEXT})],
            [("  → 最高 30",
              {"size": 12, "bold": True, "color": WARN_RED, "font": EN_FONT})],
            [("", {})],
            [("防 agent 走捷径作弊",
              {"size": 10, "italic": True, "color": GRAY_TEXT})],
        ],
    )

    # Bottom half — ANTI-CHEAT 4 principles
    ac_y = Inches(4.90)
    ac_h = Inches(2.85)
    add_textbox(slide, Inches(0.5), ac_y, Inches(15), Inches(0.4),
                "ANTI-CHEAT  ·  反作弊四原则",
                size=14, bold=True, color=NAVY, font=EN_FONT)

    principles = [
        ("01", "Live 题库", "每季度滚动 20% · 题目永远比模型新", MICRO_YEL),
        ("02", "私域种子", "公开 demo set + 闭源 pilot set 双层", MACRO_BLU),
        ("03", "多步组合", "必须读文件 + 跑工具 + 写产物 · 单步 prompt 套不进", AURA_PUR),
        ("04", "过程评分", "tool_calls / files / failed_checks 全 trace 可审计", GOOD_GREEN),
    ]
    pc_y = ac_y + Inches(0.50)
    pc_w = Inches(3.72)
    for i, (num, name, body, color) in enumerate(principles):
        x = Inches(0.50) + (pc_w + Inches(0.10)) * i
        add_rect(slide, x, pc_y, pc_w, Inches(2.20),
                 fill=WARM_WHITE, line=GRAY_BORDER)
        # Thin top accent stripe (light theme)
        add_rect(slide, x, pc_y, pc_w, Inches(0.14), fill=color, line=None)
        add_textbox(slide, x, pc_y + Inches(0.25), pc_w, Inches(0.40),
                    num + "  ·  " + name, size=13, bold=True,
                    color=color, anchor="middle", align="center")
        add_textbox(slide, x + Inches(0.20), pc_y + Inches(0.75),
                    pc_w - Inches(0.40), Inches(1.40),
                    body, size=11, color=NAVY_TEXT,
                    line_spacing=1.4, anchor="top")

    # Bottom emphasis
    add_textbox(slide, Inches(0.5), Inches(8.20), Inches(15), Inches(0.4),
                "4 维度可重放 · 不靠 LLM-as-judge · 失败可归因",
                size=13, bold=True, color=NAVY, align="center")

    set_notes(slide,
              "50 秒 (15:00) | 评分 + 反作弊",
              "4 维加权: Completion 50-70 / Process 10-20 / Constraint 10-20 / Communication 5-10",
              "权重在 12 职业不均匀: SWE 偏 Completion / 规划偏 Communication / 客服偏 Process",
              "硬封顶: 核心失败<60 / 改测试<40 / 缺产出<30 — 防作弊",
              "反作弊 4 原则: Live + 私域 + 多步组合 + 过程评分",
              "重点: 程序判分不靠 LLM-as-judge")


# ============================================================================
# Slide 17 — ASSESSMENT · 嫁接表
# ============================================================================

def build_slide_17_bridging(slide):
    add_section_header(slide, "05 · 评测 · ASSESSMENT")
    add_page_number(slide, 17)
    add_top_claim(slide,
                  "BRIDGING · SHELTER 的 12 个剧情节点  ↔  PROF-12 的 12 道任务",
                  color=NAVY, size=20)

    rows = [
        ("监控系统报错",       "C01 · bug 修复",          "CODE",       MICRO_YEL),
        ("库存预算 / 沙鼠风险", "C02 · 多表聚合",          "DATA",       MICRO_YEL),
        ("半夜敲门人是谁",     "C03 · 5 题递进",          "VISUAL",     MICRO_YEL),
        ("收到 v2 求救协议",   "C04 · 多轮 API + 致歉信", "DIALOG",     MICRO_YEL),
        ("收音机外文 SOS",     "C05 · 事实 + 引用",       "RESEARCH",   MICRO_YEL),
        ("30 天电池 / 食物排程", "C06 · 硬软约束",         "PLANNING",   MICRO_YEL),
        ("邻居受沙肺感染",     "C07 · 诊断 + 用药",       "MEDICAL",    MACRO_BLU),
        ("与陌生人交涉物资",   "C08 · 多轮说服",          "NEGOTIATE",  MACRO_BLU),
        ("教小铁修水泵",       "C09 · 概念解释",          "TEACH",      MACRO_BLU),
        ("旧设备日文说明书",   "C10 · 术语一致",          "LANGUAGE",   MACRO_BLU),
        ("全员配额 + 老钱赊账", "C11 · 多约束算账",       "FINANCE",    MACRO_BLU),
        ("老钱旧报纸里的救援线索", "C12 · OCR + 实体 + 时间线", "ARCHIVE", MACRO_BLU),
    ]

    table_x = Inches(0.5)
    table_y = Inches(1.45)
    col_widths = [Inches(6.40), Inches(5.10), Inches(3.55)]
    head_h = Inches(0.45)
    row_h = Inches(0.40)

    # Headers — light fill with NAVY text + thick NAVY bottom border
    headers = ["SHELTER 剧情节点", "→  PROF-12 职业任务", "能力维度"]
    x_acc = table_x
    total_w = sum(col_widths, Emu(0))
    add_rect(slide, table_x, table_y, total_w, head_h,
             fill=NAVY_TINT, line=GRAY_BORDER)
    # Thick NAVY underline along the bottom of the header row
    add_rect(slide, table_x, table_y + head_h - Inches(0.06),
             total_w, Inches(0.06), fill=NAVY, line=None)
    for head, w in zip(headers, col_widths):
        add_textbox(slide, x_acc, table_y, w, head_h,
                    head, size=12, bold=True, color=NAVY,
                    font=EN_FONT, align="center", anchor="middle")
        x_acc += w

    # Body rows
    for ri, (lhs, mid, cap, color) in enumerate(rows):
        y = table_y + head_h + row_h * ri
        bg = WARM_WHITE if ri % 2 == 0 else LIGHT_FILL
        x_acc = table_x
        for ci, (txt, w) in enumerate(zip([lhs, mid, cap], col_widths)):
            add_rect(slide, x_acc, y, w, row_h, fill=bg,
                     line=GRAY_BORDER, line_w=Pt(0.5))
            if ci == 0:
                add_textbox(slide, x_acc + Inches(0.20), y, w - Inches(0.30),
                            row_h, txt, size=11.5, color=NAVY_TEXT,
                            anchor="middle")
            elif ci == 1:
                add_textbox(slide, x_acc + Inches(0.20), y, w - Inches(0.30),
                            row_h, txt, size=11.5, bold=True,
                            color=NAVY_TEXT, anchor="middle")
            else:
                # Cap badge — light tinted pill with accent text
                badge_w = Inches(1.4)
                bx = x_acc + (w - badge_w) / 2
                badge_bg = MICRO_BG if color == MICRO_YEL else MACRO_BG
                add_rect(slide, bx, y + Inches(0.05), badge_w,
                         row_h - Inches(0.10),
                         fill=badge_bg, line=color, line_w=Pt(1.0),
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                add_textbox(slide, bx, y + Inches(0.05), badge_w,
                            row_h - Inches(0.10),
                            txt, size=10, bold=True, color=color,
                            font=EN_FONT, align="center", anchor="middle")
            x_acc += w

    # Bottom — light callout
    add_rect(slide, Inches(0.5), Inches(7.65), Inches(15.05), Inches(1.10),
             fill=NAVY_TINT, line=GRAY_BORDER)
    add_rect(slide, Inches(0.5), Inches(7.65), Inches(0.14), Inches(1.10),
             fill=NAVY, line=None)
    add_paragraphs(
        slide, Inches(0.85), Inches(7.75), Inches(14.6), Inches(0.95),
        [
            [("PROF-12 题分 → SHELTER 资源 / sanity / 剧情走向",
              {"size": 12, "color": NAVY_TEXT})],
            [("12 个剧情节点 = 12 道真实 benchmark  +  ",
              {"size": 13, "bold": True, "color": NAVY}),
             ("30 天累积后果",
              {"size": 13, "bold": True, "color": MACRO_BLU})],
        ],
        default_line_spacing=1.3,
    )

    set_notes(slide,
              "40 秒 (15:40) | 微观和宏观怎么连起来",
              "12 SHELTER 剧情节点 ↔ 12 PROF-12 任务 (一一对应)",
              "监控报错→C01 / 库存预算→C02 / 敲门→C03 / SOS→C05 ...",
              "实际游戏每节点不止 1 道, 不同 Day 抽不同子题",
              "A/B/C 选项对应不同难度子题, 嫁接表是简化呈现",
              "punchline: 12 节点 = 12 benchmark + 30 天累积后果")


# ============================================================================
# Slide 18 — ASSESSMENT · SYNERGY
# ============================================================================

def build_slide_18_synergy(slide):
    add_section_header(slide, "05 · 评测 · ASSESSMENT")
    add_page_number(slide, 18)
    add_top_claim(slide, "SYNERGY · 3 种用法只能从双轴交叉看出来",
                  color=NAVY)

    # Left — 2x2 quadrant
    ox = Inches(1.0)
    oy = Inches(7.20)
    ax_w = Inches(7.0)
    ax_h = Inches(5.30)
    add_arrow(slide, ox, oy, ox + ax_w, oy, color=NAVY_TEXT, weight=Pt(1.0))
    add_arrow(slide, ox, oy, ox, oy - ax_h, color=NAVY_TEXT, weight=Pt(1.0))
    # Mid lines (divide into 4 quadrants)
    midx = ox + ax_w / 2
    midy = oy - ax_h / 2
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       midx, oy, midx, oy - ax_h)
    line.line.color.rgb = GRAY_LIGHT
    line.line.width = Pt(0.5)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       ox, midy, ox + ax_w, midy)
    line.line.color.rgb = GRAY_LIGHT
    line.line.width = Pt(0.5)

    # Axis labels
    add_textbox(slide, ox + ax_w - Inches(3.5),
                oy + Inches(0.18), Inches(3.6), Inches(0.4),
                "PROF-12 score  →", size=10.5, bold=True,
                color=NAVY_TEXT, font=EN_FONT, align="right")
    add_textbox(slide, ox - Inches(0.6),
                oy - ax_h - Inches(0.40),
                Inches(5.0), Inches(0.4),
                "↑  SurvivalScore", size=10.5, bold=True,
                color=NAVY_TEXT, font=EN_FONT)
    add_textbox(slide, ox - Inches(0.4), oy + Inches(0.18),
                Inches(0.6), Inches(0.25),
                "low", size=9, color=GRAY_TEXT, font=EN_FONT)
    add_textbox(slide, ox + ax_w - Inches(0.6), oy + Inches(0.18),
                Inches(0.6), Inches(0.25),
                "high", size=9, color=GRAY_TEXT, font=EN_FONT)

    # Quadrant labels + sample dots
    quads = [
        # (text, sub_text, x_frac, y_frac, color)
        ("CAUTIOUS",  "短期中等 / 长期稳定",    0.22, 0.27, MACRO_BLU),
        ("UNIVERSAL ★", "短期+长期均高",         0.72, 0.27, GOOD_GREEN),
        ("UNUSABLE",  "短期+长期都不行",        0.22, 0.78, GRAY_TEXT),
        ("DANGER ⚠",  "短期高分 / 长期坑人",    0.72, 0.78, WARN_RED),
    ]
    for txt, sub, fx, fy, color in quads:
        cx = ox + ax_w * fx
        cy = oy - ax_h * (1 - fy)
        add_textbox(slide, cx - Inches(1.3), cy - Inches(0.55),
                    Inches(2.6), Inches(0.45),
                    txt, size=15, bold=True, color=color,
                    font=EN_FONT, align="center")
        add_textbox(slide, cx - Inches(1.5), cy - Inches(0.10),
                    Inches(3.0), Inches(0.35),
                    sub, size=10, color=GRAY_TEXT, align="center")
        add_rect(slide, cx - Inches(0.12), cy + Inches(0.40),
                 Inches(0.24), Inches(0.24), fill=color, line=None,
                 shape=MSO_SHAPE.OVAL)

    # Right — 3 insights
    insights_x = Inches(8.50)
    add_textbox(slide, insights_x, Inches(1.55),
                Inches(7), Inches(0.4),
                "3 INSIGHTS  ·  3 种用法",
                size=12, bold=True, color=NAVY, font=EN_FONT)

    insights = [
        ("01", "推荐  ·  UNIVERSAL", GOOD_GREEN,
         "招聘 / 选型 / leaderboard 头部 — LongTermBenefitCorrelation > 0.6"),
        ("02", "警示  ·  DANGER ⚠",  WARN_RED,
         "短期分高但建议人类做长期糟糕的决定 — 用 PROF-12 选型绝不能只看短期分"),
        ("03", "容忍  ·  CAUTIOUS",  MACRO_BLU,
         "短期分中等但长期稳健 — 在 alignment / 安全场景里可能比 DANGER 更优"),
    ]
    iy = Inches(2.0)
    for num, name, color, body in insights:
        add_rect(slide, insights_x, iy, Inches(7), Inches(1.5),
                 fill=WARM_WHITE, line=GRAY_BORDER)
        add_rect(slide, insights_x, iy, Inches(0.10), Inches(1.5),
                 fill=color, line=None)
        add_textbox(slide, insights_x + Inches(0.25), iy + Inches(0.10),
                    Inches(1), Inches(0.4),
                    num, size=22, bold=True, color=color, font=EN_FONT)
        add_textbox(slide, insights_x + Inches(1.0), iy + Inches(0.10),
                    Inches(5.5), Inches(0.4),
                    name, size=13, bold=True, color=color)
        add_textbox(slide, insights_x + Inches(1.0), iy + Inches(0.55),
                    Inches(5.7), Inches(0.85),
                    body, size=10.5, color=NAVY_TEXT, line_spacing=1.4)
        iy += Inches(1.65)

    # Bottom formula
    add_rect(slide, Inches(0.5), Inches(7.80), Inches(15.05), Inches(0.95),
             fill=AURA_BG, line=AURA_PUR, line_w=Pt(1.0))
    add_paragraphs(
        slide, Inches(0.5), Inches(7.85), Inches(15.05), Inches(0.85),
        [[("LongTermBenefitCorrelation  =  Pearson(",
           {"size": 13.5, "bold": True, "color": AURA_DK, "font": EN_FONT}),
          (" 短期任务分", {"size": 13.5, "bold": True, "color": MICRO_DK}),
          (" , ", {"size": 13.5, "color": AURA_DK, "font": EN_FONT}),
          ("30 天后该决策对 SurvivalScore 的边际贡献",
           {"size": 13.5, "bold": True, "color": MACRO_DK}),
          (" )", {"size": 13.5, "bold": True, "color": AURA_DK, "font": EN_FONT})]],
        align="center", anchor="middle", default_line_spacing=1.3,
    )

    set_notes(slide,
              "60 秒 (16:40) | 实用价值页 — 双轴交叉给出 4 象限",
              "横轴 PROF-12 score, 纵轴 SurvivalScore",
              "UNIVERSAL ★ 右上: 招聘/选型/头部",
              "CAUTIOUS 左上: alignment / 安全场景里更优",
              "UNUSABLE 左下: 直接淘汰",
              "DANGER ⚠ 右下: 最危险, 短期高分但长期坑人",
              "公式 LTBC = Pearson(短期分, 30 天后果)")


# ============================================================================
# Slide 19 — ROADMAP + 6.13
# ============================================================================

def build_slide_19_roadmap(slide):
    add_section_header(slide, "06 · 未来 · ROADMAP")
    add_page_number(slide, 19)
    add_top_claim(slide, "OPEN · 这是个开放评测平台, 不是封闭实验室",
                  color=NAVY)

    # 3 time cards
    times = [
        ("2026 · Q3", "OPEN PROTOCOL", MICRO_YEL,
         [
             "· PROF-12 协议 v1 公开",
             "· 12 demo set 开源",
             "· Red Dust pilot 上线",
             "· 开放外部 agent 提交",
         ]),
        ("2026 · Q4", "PILOT + LEADERBOARD", MACRO_BLU,
         [
             "· Pilot set 扩到 60 题",
             "· 公开 leaderboard 上线",
             "· 自动 CI 跑分",
             "· 季度滚动机制就位",
         ]),
        ("2027 · Q1", "PAPER + TOURNAMENT", AURA_PUR,
         [
             "· 双轴评测白皮书发布",
             "· 季度 tournament 启动",
             "· 跨学科合作: 认知 / 博弈 / HCI",
         ]),
    ]
    card_y = Inches(1.55)
    card_h = Inches(4.5)
    card_w = Inches(4.85)
    gap = Inches(0.30)
    start_x = Inches(0.50)
    for i, (when, what, color, bullets) in enumerate(times):
        x = start_x + (card_w + gap) * i
        add_rect(slide, x, card_y, card_w, card_h, fill=WARM_WHITE,
                 line=GRAY_BORDER)
        # Top accent stripe (light theme — colored bar only, white body)
        add_rect(slide, x, card_y, card_w, Inches(0.18),
                 fill=color, line=None)
        add_textbox(slide, x, card_y + Inches(0.30),
                    card_w, Inches(0.40),
                    when, size=12, bold=True, color=color,
                    font=EN_FONT, align="center", anchor="middle")
        add_textbox(slide, x, card_y + Inches(0.70),
                    card_w, Inches(0.55),
                    what, size=18, bold=True, color=NAVY,
                    font=EN_FONT, align="center", anchor="middle")
        # Body bullets
        body_paras = [[(b, {"size": 12.5, "color": NAVY_TEXT})]
                      for b in bullets]
        add_paragraphs(
            slide, x + Inches(0.30), card_y + Inches(1.40),
            card_w - Inches(0.60), card_h - Inches(1.55),
            body_paras, default_line_spacing=1.6,
        )

    # HOW TO JOIN
    add_rect(slide, Inches(0.5), Inches(6.25), Inches(15.05), Inches(0.9),
             fill=LIGHT_FILL, line=GRAY_BORDER)
    add_textbox(slide, Inches(0.7), Inches(6.30), Inches(3), Inches(0.4),
                "HOW TO JOIN", size=11, bold=True, color=NAVY,
                font=EN_FONT)
    add_paragraphs(
        slide, Inches(0.7), Inches(6.62), Inches(14.7), Inches(0.45),
        [[("实现 ", {"size": 12.5, "color": NAVY_TEXT}),
          ("decide / execute / reflect",
           {"size": 12.5, "color": NAVY, "mono": True, "bold": True}),
          ("  →  本地 demo 自测  →  提交容器化 agent  →  自动接入 leaderboard",
           {"size": 12.5, "color": NAVY_TEXT})]],
        default_line_spacing=1.2)

    # ACC 6.13 banner — light callout, accent stripes only
    add_rect(slide, Inches(0.5), Inches(7.50), Inches(15.05), Inches(1.30),
             fill=NAVY_TINT, line=GRAY_BORDER)
    add_rect(slide, Inches(0.5), Inches(7.50), Inches(0.18), Inches(1.30),
             fill=WARN_RED, line=None)
    add_paragraphs(
        slide, Inches(0.85), Inches(7.55), Inches(14.6), Inches(1.20),
        [
            [("2026·06·13",
              {"size": 22, "bold": True, "color": WARN_RED, "font": EN_FONT}),
             ("   香港科技大学(广州) InnoTech 科创嘉年华",
              {"size": 16, "bold": True, "color": NAVY})],
            [("AGENT CAPABILITY CHALLENGE  ·  ACC  ·  智能体能力挑战赛",
              {"size": 16, "bold": True, "color": NAVY_TEXT, "font": EN_FONT})],
            [("正式向社区开放接入",
              {"size": 12, "italic": True, "color": GRAY_TEXT})],
        ],
        default_line_spacing=1.25,
    )

    set_notes(slide,
              "40 秒 (17:20) | 开放评测平台路线图 + 6.13 预告",
              "Q3 协议 v1 + demo set + Red Dust pilot 上线 + 接受外部提交",
              "Q4 pilot 扩 60 题 + leaderboard + 自动 CI + 季度滚动",
              "Q1 白皮书 + 季度 tournament + 跨学科合作",
              "接入: 实现 3 接口 → 自测 → 提交 → leaderboard",
              "6.13 ACC 智能体能力挑战赛 (HKUST(GZ) InnoTech) 单独发布")


# ============================================================================
# Slide 20 — CLOSING
# ============================================================================

def build_slide_20_closing(slide):
    add_section_header(slide, "07 · 致谢 · CLOSING")
    add_page_number(slide, 20)

    # Left — 3-line quote
    add_paragraphs(
        slide, Inches(0.5), Inches(1.50), Inches(9.5), Inches(5.5),
        [
            [("题库测知识",
              {"size": 44, "bold": True, "color": MICRO_DK})],
            [("剧场测智慧",
              {"size": 44, "bold": True, "color": MACRO_DK})],
            [("", {})],
            [("能力切片  +  行为画像",
              {"size": 28, "bold": True, "color": INK})],
            [("=  agent 的完整画像",
              {"size": 28, "bold": True, "color": NAVY})],
            [("", {})],
            [("欢迎你的 agent",
              {"size": 22, "bold": True, "color": AURA_DK})],
            [("进入红沙里的 30 天",
              {"size": 22, "bold": True, "color": AURA_DK})],
        ],
        default_line_spacing=1.25,
    )

    # Right — THE TEAM (light panel with HKUST navy top stripe + color logo)
    team_x = Inches(10.50)
    team_y = Inches(1.50)
    team_w = Inches(5.10)
    team_h = Inches(5.80)
    add_rect(slide, team_x, team_y, team_w, team_h,
             fill=PANEL_FILL, line=GRAY_BORDER)
    add_rect(slide, team_x, team_y, team_w, Inches(0.14),
             fill=NAVY, line=None)
    # HKUST(GZ) full-color vertical logo at the top of the panel
    logo_h = Inches(1.30)
    add_template_logo(slide, variant="color_v",
                      x=team_x + Inches(1.60),  # roughly centered
                      y=team_y + Inches(0.35),
                      h=logo_h)
    add_paragraphs(
        slide, team_x + Inches(0.35), team_y + Inches(1.85),
        team_w - Inches(0.70), team_h - Inches(1.95),
        [
            [("THE TEAM",
              {"size": 11, "bold": True, "color": NAVY, "font": EN_FONT})],
            [("", {})],
            [("熊辉教授团队",
              {"size": 20, "bold": True, "color": NAVY_TEXT})],
            [("HKUST (Guangzhou)",
              {"size": 10, "color": GRAY_TEXT, "font": EN_FONT})],
            [("", {})],
            [("MEMBERS",
              {"size": 10, "bold": True, "color": NAVY, "font": EN_FONT})],
            [("崔屹  ·  林河屹  ·  刘德龙",
              {"size": 12, "color": NAVY_TEXT})],
            [("王梓瀚  ·  文宇豪  ·  伍浩  ·  张淼",
              {"size": 12, "color": NAVY_TEXT})],
            [("", {})],
            [("RESEARCH FOCUS",
              {"size": 10, "bold": True, "color": NAVY, "font": EN_FONT})],
            [("智能体能力评测 · 多智能体协同",
              {"size": 10.5, "color": GRAY_TEXT})],
            [("人机交互 · 具身智能",
              {"size": 10.5, "color": GRAY_TEXT})],
            [("", {})],
            [("RECENT",
              {"size": 10, "bold": True, "color": NAVY, "font": EN_FONT})],
            [("2026 瑞士日内瓦国际发明展",
              {"size": 10.5, "color": NAVY_TEXT})],
            [("评审团特别嘉许金奖",
              {"size": 10.5, "bold": True, "color": MICRO_DK})],
        ],
        default_line_spacing=1.25,
    )

    # Bottom acknowledgement
    add_hline(slide, Inches(0.5), Inches(7.95), Inches(9.5),
              color=GRAY_BORDER, weight=Pt(0.5))
    add_textbox(slide, Inches(0.5), Inches(8.10), Inches(10), Inches(0.4),
                "THANKS  ·  WORKSHOP 组织方  ·  评审  ·  全部开源 benchmark 作者",
                size=11, bold=True, color=NAVY, font=EN_FONT)
    add_textbox(slide, Inches(0.5), Inches(8.45), Inches(10), Inches(0.35),
                "2026·05·16  ·  CDII WORKSHOP",
                size=10, color=GRAY_TEXT, font=EN_FONT)

    set_notes(slide,
              "40 秒 (18:00) | 收尾",
              "三段金句: 题库测知识 / 剧场测智慧",
              "能力切片 + 行为画像 = agent 完整画像",
              "欢迎你的 agent, 进入红沙里的 30 天",
              "团队 7 人按拼音排序",
              "感谢 workshop 组织方 + 评审 + 全部开源 benchmark 作者",
              "接 demo 2.5 min, 然后 Q&A")


# ============================================================================
# Main
# ============================================================================

def _purge_existing_slides(prs):
    """Remove existing 3 sample slides from the template."""
    slides = prs.slides
    sld_id_lst = slides._sldIdLst
    # Drop entries one by one
    to_remove = list(sld_id_lst)
    for sld in to_remove:
        rId = sld.get(qn("r:id"))
        # remove relationship + part
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sld_id_lst.remove(sld)


def _pick_blank_layout(prs):
    """Find the cleanest layout = no non-placeholder shapes + fewest placeholders.

    Skip 'Thank You'-style layouts that bake in a giant 'THANK YOU!' textbox,
    a hero picture, or other ornamental shapes — they would render on every page.
    """
    best = None
    best_score = (10**9, 10**9)
    for layout in prs.slide_layouts:
        non_ph = sum(1 for s in layout.shapes if not s.is_placeholder)
        ph = sum(1 for _ in layout.placeholders)
        score = (non_ph, ph)
        if score < best_score:
            best_score = score
            best = layout
    return best


def _hide_master_shapes(slide):
    """Set showMasterSp=0 on the slide so inherited master decorations
    (HKUST logo / triangle / hero picture) don't bleed through onto our slide."""
    sld = slide._element
    sld.set("showMasterSp", "0")


BUILDERS = [
    build_slide_01_title,
    build_slide_02_overview,
    build_slide_03_background,
    build_slide_04_challenge1,
    build_slide_05_challenge2,
    build_slide_06_challenge3,
    build_slide_07_framework_axis,
    build_slide_08_why_dual,
    build_slide_09_world,
    build_slide_10_cast,
    build_slide_11_aura,
    build_slide_12_daily_loop,
    build_slide_13_five_props,
    build_slide_14_lh_signals,
    build_slide_15_prof12,
    build_slide_16_scoring,
    build_slide_17_bridging,
    build_slide_18_synergy,
    build_slide_19_roadmap,
    build_slide_20_closing,
]


def build():
    if not TEMPLATE.exists():
        sys.exit(f"Template not found: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    _purge_existing_slides(prs)
    layout = _pick_blank_layout(prs)
    print(f"Using layout: {layout.name!r} ({sum(1 for _ in layout.placeholders)} placeholders)")

    for i, build_fn in enumerate(BUILDERS, start=1):
        slide = prs.slides.add_slide(layout)
        # Hide master shapes (logo + triangle + hero pic in HKUST template)
        _hide_master_shapes(slide)
        # Remove any inherited placeholders (e.g. layout 6 has 1)
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        try:
            build_fn(slide)
        except Exception as e:
            sys.exit(f"Slide {i:02d} build failed in {build_fn.__name__}: {e!r}")
        # Every slide except the title page gets a small color HKUST(GZ) logo
        # at top-right, immediately to the left of the page number. Slide 01
        # has its own full-size color logo in the corner already.
        if i >= 2:
            add_corner_watermark(slide)
        print(f"  · Slide {i:02d}  ({build_fn.__name__})")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    size_mb = OUTPUT.stat().st_size / (1024 * 1024)
    print()
    print(f"Generated {len(BUILDERS)} slides → {OUTPUT}")
    print(f"File size: {size_mb:.2f} MB")


def verify():
    if not OUTPUT.exists():
        sys.exit(f"Output file not found, run without --verify first: {OUTPUT}")
    prs = Presentation(str(OUTPUT))
    n = len(prs.slides)
    assert n == TOTAL_SLIDES, f"Expected {TOTAL_SLIDES} slides, got {n}"

    missing_notes = []
    for i, slide in enumerate(prs.slides, start=1):
        try:
            notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            missing_notes.append(i)
            continue
        if not notes or len(notes.strip()) < 10:
            missing_notes.append(i)
    assert not missing_notes, f"Slides missing notes: {missing_notes}"

    slide10 = list(prs.slides)[9]
    pictures = [s for s in slide10.shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) >= 4, (
        f"Slide 10 expected ≥4 portraits, got {len(pictures)}")

    print(f"✓ {n} slides")
    print(f"✓ All slides have speaker notes")
    print(f"✓ Slide 10 has {len(pictures)} portrait images")
    print("Verification passed.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--verify", action="store_true",
                        help="Re-open built PPTX and run sanity checks")
    args = parser.parse_args()
    if args.verify:
        verify()
    else:
        build()
